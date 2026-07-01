"""Build a clean presentation on how the X2 walking policy was trained.

Plain, correct wording. Consistent layout with real bullets and hanging indents.
Math rendered with matplotlib. Progression and parallel thumbnails pulled from the
recorded videos. Output: presentation/X2_Walking_Policy.pptx
"""

import matplotlib

matplotlib.use("Agg")
from pathlib import Path

import matplotlib.pyplot as plt
import mediapy as media
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

VIDEO_DIR = Path("/home/baaqer/baaqer_ws/policy_progression_videos")
OUT_DIR = Path("/home/baaqer/baaqer_ws/presentation")
THUMB_DIR = OUT_DIR / "thumbs"
EQ_DIR = OUT_DIR / "equations"

# Palette.
NAVY = RGBColor(0x15, 0x22, 0x40)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
DARK = RGBColor(0x23, 0x28, 0x2E)
GRAY = RGBColor(0x5C, 0x66, 0x70)
LIGHT = RGBColor(0xEF, 0xF3, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"
EQ_HEX = "#20242B"

IN = 914400
SLIDE_W = Emu(int(13.333 * IN))
SLIDE_H = Emu(int(7.5 * IN))


def emu(inches):
  return Emu(int(inches * IN))


# ----------------------------- assets -----------------------------
def extract_thumbnails():
  THUMB_DIR.mkdir(parents=True, exist_ok=True)
  clips = [
    ("01_iter0_failing", "thumb_0", 140),
    ("02_iter100", "thumb_100", 140),
    ("03_iter3000", "thumb_3000", 140),
    ("04_iter6000", "thumb_6000", 140),
    ("05_iter9999_final", "thumb_9999", 140),
    ("06_parallel_training", "thumb_parallel", 120),
  ]
  paths = {}
  for src, name, idx in clips:
    v = media.read_video(str(VIDEO_DIR / f"{src}.mp4"))
    frame = np.asarray(v[min(idx, len(v) - 1)])
    rows = np.where(frame.reshape(frame.shape[0], -1).sum(axis=1) > 0)[0]
    if len(rows) > 0:
      frame = frame[rows.min():rows.max() + 1]
    out = THUMB_DIR / f"{name}.png"
    Image.fromarray(frame).save(out)
    paths[name] = out
  return paths


def render_equations():
  EQ_DIR.mkdir(parents=True, exist_ok=True)
  eqs = {
    "policy": r"$a_t \sim \pi_\theta(\,\cdot \mid s_t\,)$",
    "targets": r"$q^{\mathrm{target}} = q^{\mathrm{default}} + \alpha \odot a_t$",
    "reward": r"$r_t = \sum_k w_k \, r_k(s_t, a_t)$",
    "velterm": r"$r_{\mathrm{vel}} = \exp\!\left(-\,\frac{\| v^{*} - v \|^{2}}{\sigma^{2}}\right)$",
    "return": r"$J(\theta) = \mathbb{E}\!\left[\, \sum_t \gamma^{\,t} \, r_t \,\right]$",
    "ppo": r"$L(\theta) = \mathbb{E}\left[\, \min(\rho_t \hat{A}_t,\ \mathrm{clip}(\rho_t, 1-\epsilon, 1+\epsilon)\,\hat{A}_t) \,\right]$",
    "ratio": r"$\rho_t = \frac{\pi_\theta(a_t \mid s_t)}{\pi_{\theta_{\mathrm{old}}}(a_t \mid s_t)}$",
  }
  paths = {}
  for name, tex in eqs.items():
    fig = plt.figure(figsize=(8, 1.2))
    fig.text(0.0, 0.5, tex, va="center", ha="left", fontsize=30, color=EQ_HEX)
    out = EQ_DIR / f"{name}.png"
    fig.savefig(out, dpi=220, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    paths[name] = out
  return paths


# ----------------------------- primitives -----------------------------
def add_bg(slide, color):
  s = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
  s.fill.solid()
  s.fill.fore_color.rgb = color
  s.line.fill.background()
  s.shadow.inherit = False
  slide.shapes._spTree.remove(s._element)
  slide.shapes._spTree.insert(2, s._element)
  return s


def add_rect(slide, x, y, w, h, color):
  s = slide.shapes.add_shape(1, x, y, w, h)
  s.fill.solid()
  s.fill.fore_color.rgb = color
  s.line.fill.background()
  s.shadow.inherit = False
  return s


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=8, line_spacing=1.05):
  tb = slide.shapes.add_textbox(x, y, w, h)
  tf = tb.text_frame
  tf.word_wrap = True
  tf.vertical_anchor = anchor
  for i, para in enumerate(runs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    p.line_spacing = line_spacing
    for (text, size, bold, color) in para:
      r = p.add_run()
      r.text = text
      r.font.size = Pt(size)
      r.font.bold = bold
      r.font.name = FONT
      r.font.color.rgb = color
  return tb


def title_band(slide, title):
  add_rect(slide, 0, 0, SLIDE_W, emu(1.1), NAVY)
  add_rect(slide, 0, emu(1.1), SLIDE_W, emu(0.055), ACCENT)
  add_text(slide, emu(0.7), emu(0.24), emu(12), emu(0.62),
           [[(title, 29, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n):
  add_text(slide, emu(0.7), emu(7.04), emu(9), emu(0.36),
           [[("X2 humanoid walking policy", 9, False, GRAY)]])
  add_text(slide, emu(11.6), emu(7.04), emu(1.0), emu(0.36),
           [[(str(n), 9, False, GRAY)]], align=PP_ALIGN.RIGHT)


def _apply_bullet(p, level):
  pPr = p._p.get_or_add_pPr()
  pPr.set("marL", str(int((0.32 + level * 0.42) * IN)))
  pPr.set("indent", str(-int(0.32 * IN)))
  buClr = pPr.makeelement(qn("a:buClr"), {})
  srgb = pPr.makeelement(qn("a:srgbClr"), {"val": "2E86C1" if level == 0 else "8A97A6"})
  buClr.append(srgb)
  buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
  buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022" if level == 0 else "\u25AA"})
  pPr.append(buClr)
  pPr.append(buFont)
  pPr.append(buChar)


def bullets(slide, x, y, w, h, items):
  """items: list of (level, text, size, bold, color)."""
  tb = slide.shapes.add_textbox(x, y, w, h)
  tf = tb.text_frame
  tf.word_wrap = True
  for i, (level, text, size, bold, color) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6 if level else 12)
    p.space_before = Pt(2 if level else 6)
    p.line_spacing = 1.08
    _apply_bullet(p, level)
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
  return tb


def bullet_slide(prs, title, items, n):
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  add_bg(slide, WHITE)
  title_band(slide, title)
  bullets(slide, emu(0.85), emu(1.55), emu(11.6), emu(5.2), items)
  footer(slide, n)
  return slide


def place_image_fit(slide, path, x, y, max_w, max_h):
  img = Image.open(path)
  ar = img.width / img.height
  w = max_w
  h = w / ar
  if h > max_h:
    h = max_h
    w = h * ar
  x_c = x + (max_w - w) / 2
  slide.shapes.add_picture(str(path), emu(x_c), emu(y), width=emu(w), height=emu(h))
  return w, h


# ----------------------------- deck -----------------------------
def main():
  OUT_DIR.mkdir(parents=True, exist_ok=True)
  thumbs = extract_thumbnails()
  eqs = render_equations()

  prs = Presentation()
  prs.slide_width = SLIDE_W
  prs.slide_height = SLIDE_H

  B = DARK  # body color
  G = GRAY  # sub color

  # ---- Slide 1: title ----
  s = prs.slides.add_slide(prs.slide_layouts[6])
  add_bg(s, NAVY)
  add_rect(s, emu(0.9), emu(2.35), emu(2.2), emu(0.08), ACCENT)
  add_text(s, emu(0.9), emu(2.55), emu(11.6), emu(1.5),
           [[("Teaching the X2 Humanoid to Walk", 42, True, WHITE)]])
  add_text(s, emu(0.92), emu(3.85), emu(11), emu(0.9),
           [[("Training a walking controller in simulation and running it on the "
              "real robot", 19, False, RGBColor(0xC6, 0xD2, 0xE2))]])
  add_rect(s, emu(0.92), emu(5.35), emu(4.4), emu(0.02), RGBColor(0x3A, 0x4A, 0x66))
  add_text(s, emu(0.92), emu(5.5), emu(11), emu(0.5),
           [[("Baaqer Farhat", 20, True, WHITE), ("      Mahdi Taheri", 20, True, WHITE)]])
  add_text(s, emu(0.92), emu(6.15), emu(11), emu(0.45),
           [[("AgiBot X2      MuJoCo      mjlab      Reinforcement learning",
              13, False, ACCENT)]])

  # ---- Slide 2: objective ----
  bullet_slide(prs, "Objective", [
    (0, "Learn a velocity conditioned walking controller for the 31 DoF X2 humanoid.", 19, False, B),
    (0, "Track a commanded forward, lateral, and yaw velocity while keeping the base stable.", 19, False, B),
    (0, "Train end to end in simulation with reinforcement learning, then transfer to hardware.", 19, False, B),
    (0, "The gait emerges from the reward objective. There are no reference trajectories or scripted state machines.", 19, False, B),
  ], 2)

  # ---- Slide 3: control policy ----
  bullet_slide(prs, "The control policy", [
    (0, "A feedforward network mapping proprioceptive state to joint position targets at 50 Hz.", 19, False, B),
    (0, "Observation:", 19, True, B),
    (1, "Base angular velocity and projected gravity.", 17, False, G),
    (1, "Joint positions and velocities, and the previous action.", 17, False, G),
    (1, "The velocity command driving the gait.", 17, False, G),
    (0, "Action:", 19, True, B),
    (1, "A position target per joint, tracked by the actuator PD loops.", 17, False, G),
    (0, "Actor and critic share the observation. The critic additionally sees base linear velocity as a privileged, training only signal.", 17, False, G),
  ], 3)

  # ---- Slide 4: parallel setup (with swarm image) ----
  s = prs.slides.add_slide(prs.slide_layouts[6])
  add_bg(s, WHITE)
  title_band(s, "Simulation and training setup")
  bullets(s, emu(0.85), emu(1.55), emu(6.1), emu(5.0), [
    (0, "Training runs in MuJoCo Warp through mjlab, entirely on a single GPU.", 18, False, B),
    (0, "Roughly 4096 environments step in parallel, each an independent robot on its own terrain.", 18, False, B),
    (0, "One shared actor and critic collect experience from all environments at once.", 18, False, B),
    (0, "This yields billions of simulated timesteps within a few hours of wall clock training.", 18, False, B),
  ])
  w, h = place_image_fit(s, thumbs["thumb_parallel"], 7.05, 1.7, 5.9, 3.6)
  add_text(s, emu(7.05), emu(1.7 + h + 0.12), emu(5.9), emu(0.5),
           [[("A subset of the parallel environments stepping simultaneously.",
              12, False, GRAY)]], align=PP_ALIGN.CENTER)
  footer(s, 4)

  # ---- Slide 5: learning formulation (math + PPO) ----
  s = prs.slides.add_slide(prs.slide_layouts[6])
  add_bg(s, WHITE)
  title_band(s, "Learning formulation")
  rows = [
    ("Stochastic policy over joint targets", "policy"),
    ("Actions define PD position setpoints", "targets"),
    ("Reward is a weighted sum of shaped terms", "reward"),
    ("Velocity tracking via an exponential kernel", "velterm"),
    ("PPO maximizes the clipped surrogate", "ppo"),
  ]
  y = 1.5
  row_h = 0.9
  for label, key in rows:
    add_rect(s, emu(0.85), emu(y + 0.05), emu(0.09), emu(row_h - 0.28), ACCENT)
    add_text(s, emu(1.1), emu(y), emu(5.0), emu(row_h - 0.18),
             [[(label, 14, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
    place_image_fit(s, eqs[key], 6.1, y + 0.08, 6.8, row_h - 0.30)
    y += row_h
  add_rect(s, emu(0.85), emu(6.15), emu(11.6), emu(0.02), RGBColor(0xD5, 0xDD, 0xE5))
  add_text(s, emu(0.85), emu(6.28), emu(11.7), emu(0.7),
           [[("On policy optimization: rollouts are batched across all environments, "
              "advantages estimated with GAE, and the policy updated by PPO with a clipped "
              "surrogate and entropy regularization.", 13, False, GRAY)]],
           line_spacing=1.1)
  footer(s, 5)

  # ---- Slide 6: reward design ----
  s = prs.slides.add_slide(prs.slide_layouts[6])
  add_bg(s, WHITE)
  title_band(s, "Reward design")
  data = [
    ("Reward term", "Purpose"),
    ("Linear and angular velocity tracking", "Match the commanded base velocity"),
    ("Base orientation and height", "Penalize torso tilt and unwanted vertical motion"),
    ("Foot clearance and air time", "Enforce distinct swing and stance phases"),
    ("Action rate and torque", "Regularize toward smooth, efficient motion"),
    ("Foot slip and contact impact", "Penalize sliding and hard landings"),
    ("Termination on fall", "End and penalize the episode when the base drops"),
  ]
  gt = s.shapes.add_table(len(data), 2, emu(0.85), emu(1.55),
                          emu(11.6), emu(4.75)).table
  gt.columns[0].width = emu(4.7)
  gt.columns[1].width = emu(6.9)
  for r in range(len(data)):
    for c in range(2):
      cell = gt.cell(r, c)
      cell.text = data[r][c]
      para = cell.text_frame.paragraphs[0]
      para.font.name = FONT
      para.font.size = Pt(15 if r else 16)
      para.font.bold = (r == 0) or (c == 0)
      if r == 0:
        para.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
      else:
        para.font.color.rgb = DARK if c == 1 else ACCENT
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
      cell.vertical_anchor = MSO_ANCHOR.MIDDLE
      cell.margin_left = emu(0.16)
      cell.margin_top = emu(0.04)
      cell.margin_bottom = emu(0.04)
  footer(s, 6)

  # ---- Slide 7: robustness ----
  bullet_slide(prs, "Robustness and transfer", [
    (0, "Velocity curriculum: the command range widens as tracking performance improves.", 19, False, B),
    (0, "Domain randomization: ground friction, added base mass, and actuator gains vary per environment.", 19, False, B),
    (0, "External perturbations: random base pushes force the policy to learn balance recovery.", 19, False, B),
    (0, "Together these reduce the simulation to reality gap and prevent overfitting to one exact setup.", 19, False, B),
  ], 7)

  # ---- Slide 8: command ----
  bullet_slide(prs, "Command conditioned locomotion", [
    (0, "The policy is conditioned on a velocity command: forward, lateral, and yaw rate.", 19, False, B),
    (0, "A single network produces forward, lateral, and turning gaits by varying that command.", 19, False, B),
    (0, "At deployment we stream a fixed command, for example 0.3 m/s forward, to drive walking.", 19, False, B),
  ], 8)

  # ---- Slide 9: progression ----
  s = prs.slides.add_slide(prs.slide_layouts[6])
  add_bg(s, WHITE)
  title_band(s, "Policy progression over training")
  items = [
    ("thumb_0", "Iteration 0", "Untrained, collapses immediately"),
    ("thumb_100", "Iteration 100", "Unstable, fails to hold balance"),
    ("thumb_3000", "Iteration 3000", "Emerging forward stepping"),
    ("thumb_6000", "Iteration 6000", "Consistent locomotion"),
    ("thumb_9999", "Iteration 10000", "Stable, converged gait"),
  ]
  n = len(items)
  margin = 0.55
  gap = 0.22
  cw = (13.333 - 2 * margin - (n - 1) * gap) / n
  top = 1.9
  for i, (key, label, sub) in enumerate(items):
    x = margin + i * (cw + gap)
    img = Image.open(thumbs[key])
    ar = img.height / img.width
    ph = cw * ar
    s.shapes.add_picture(str(thumbs[key]), emu(x), emu(top), width=emu(cw))
    ty = top + ph + 0.16
    add_text(s, emu(x), emu(ty), emu(cw), emu(0.35),
             [[(label, 14, True, NAVY)]], align=PP_ALIGN.CENTER)
    add_text(s, emu(x), emu(ty + 0.36), emu(cw), emu(0.7),
             [[(sub, 12, False, GRAY)]], align=PP_ALIGN.CENTER)
  add_text(s, emu(0.7), emu(6.45), emu(12), emu(0.5),
           [[("Deterministic rollouts of saved checkpoints under a fixed forward command.",
              13, False, GRAY)]])
  footer(s, 9)

  # ---- Slide 10: deployment + video placeholder ----
  s = prs.slides.add_slide(prs.slide_layouts[6])
  add_bg(s, WHITE)
  title_band(s, "Deployment on hardware")
  bullets(s, emu(0.85), emu(1.55), emu(6.0), emu(5.0), [
    (0, "The trained actor is exported to a compact model and runs with numpy only on the robot.", 18, False, B),
    (0, "The deployment policy drops base linear velocity from its observation, since it is not measurable on hardware. The critic used it only in training.", 18, False, B),
    (0, "Inference runs at 50 Hz over ROS 2, and the onboard PD loops track the joint targets.", 18, False, B),
    (0, "The velocity command is streamed live, so the same policy walks in any direction on the real robot.", 18, False, B),
  ])
  # Video placeholder.
  px, py, pw, ph = 7.05, 1.9, 5.9, 3.32
  box = add_rect(s, emu(px), emu(py), emu(pw), emu(ph), LIGHT)
  box.line.fill.solid()
  box.line.color.rgb = ACCENT
  box.line.width = Pt(1.75)
  add_text(s, emu(px), emu(py), emu(pw), emu(ph),
           [[("\u25B6", 30, False, ACCENT)],
            [("Hardware demo video", 16, True, NAVY)],
            [("insert clip here", 12, False, GRAY)]],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=6)
  add_text(s, emu(px), emu(py + ph + 0.12), emu(pw), emu(0.5),
           [[("Policy running on the physical X2.", 12, False, GRAY)]],
           align=PP_ALIGN.CENTER)
  footer(s, 10)

  out = OUT_DIR / "X2_Walking_Policy.pptx"
  prs.save(str(out))
  print(f"saved {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
  main()
