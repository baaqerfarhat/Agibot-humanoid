#!/usr/bin/env python3
"""Build the SURF final presentation.

    python3 presentation/surf_figures.py       # charts + equation images
    python3 presentation/build_surf_slides.py  # -> presentation/SURF_Final_Baaqer_Farhat.pptx

Video slots show a real frame from the matching rollout as a poster, with the
source file named underneath. presentation/insert_surf_videos.py swaps the
actual .mp4 files into those slots.
"""

from __future__ import annotations

import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
ASSETS = HERE / "surf_assets"
POSTERS = ASSETS / "posters"
VIDEOS = ROOT / "box_pickup" / "videos"
CLIPS = ASSETS / "clips"

# insert_surf_videos.py sets these to embed the real .mp4 files in the slots
EMBED = os.environ.get("SURF_EMBED_VIDEOS") == "1"
OUT = HERE / (
    "SURF_Final_Baaqer_Farhat_with_videos.pptx" if EMBED
    else "SURF_Final_Baaqer_Farhat.pptx"
)

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x15, 0x22, 0x40)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
GREEN = RGBColor(0x1E, 0x84, 0x49)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xB9, 0x77, 0x0E)
GRAY = RGBColor(0x5C, 0x66, 0x70)
RULE = RGBColor(0xC5, 0xCD, 0xD6)
LIGHT = RGBColor(0xEF, 0xF3, 0xF7)
PALE_B = RGBColor(0xDD, 0xEB, 0xF6)
PALE_G = RGBColor(0xEA, 0xF5, 0xEC)
PALE_R = RGBColor(0xFB, 0xEB, 0xE9)
PALE_A = RGBColor(0xFB, 0xF1, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x11, 0x14, 0x18)

FONT = "Arial"
MONO = "Consolas"

# ---------------------------------------------------------------- geometry
IN = 914400
SW, SH = 13.333, 7.5
ML, MR = 0.62, 0.62
CW = SW - ML - MR          # 12.093
BODY_TOP = 1.44
FOOTER_Y = 6.98


def emu(v: float) -> Emu:
    return Emu(int(round(v * IN)))


# ---------------------------------------------------------------- helpers
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = emu(SW)
    prs.slide_height = emu(SH)
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text="", size=14, color=NAVY, bold=False, italic=False, first=False,
         space_before=0, space_after=0, align=None, font=FONT, line=1.28):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if align is not None:
        p.alignment = align
    if text:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return p


def rich(tf, chunks, size=14, first=False, space_before=0, space_after=0,
         align=None, line=1.28):
    """chunks: list of (text, {bold, color, italic, size, font})."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if align is not None:
        p.alignment = align
    for text, style in chunks:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(style.get("size", size))
        r.font.bold = style.get("bold", False)
        r.font.italic = style.get("italic", False)
        r.font.color.rgb = style.get("color", NAVY)
        r.font.name = style.get("font", FONT)
    return p


def shape(slide, kind, x, y, w, h, fill=WHITE, line=None, lw=1.25, radius=None):
    sp = slide.shapes.add_shape(kind, emu(x), emu(y), emu(w), emu(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = emu(0.08)
    tf.margin_top = tf.margin_bottom = emu(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return sp


def rect(slide, x, y, w, h, **kw):
    return shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, **kw)


def rrect(slide, x, y, w, h, radius=0.09, **kw):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, radius=radius, **kw)


def arrow(slide, x1, y1, x2, y2, color=ACCENT, lw=1.6, dashed=False):
    cx = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    emu(x1), emu(y1), emu(x2), emu(y2))
    cx.line.color.rgb = color
    cx.line.width = Pt(lw)
    ln = cx.line._get_or_add_ln()
    if dashed:
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(d)
    head = ln.makeelement(qn("a:tailEnd"),
                          {"type": "triangle", "w": "med", "len": "med"})
    ln.append(head)
    return cx


def _fit(path, x, y, w, h):
    """Rect for an image of path's aspect ratio, centred inside (x, y, w, h)."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    bw, bh = (w, w / ar) if (w / ar) <= h else (h * ar, h)
    return x + (w - bw) / 2, y + (h - bh) / 2, bw, bh


def picture(slide, path, x, y, w=None, h=None):
    """Place an image fitted inside (w, h), centred, preserving aspect."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w is not None and h is not None:
        bx, by, bw, bh = _fit(path, x, y, w, h)
        return slide.shapes.add_picture(str(path), emu(bx), emu(by), emu(bw), emu(bh))
    if w is not None:
        return slide.shapes.add_picture(str(path), emu(x), emu(y), emu(w), emu(w / ar))
    return slide.shapes.add_picture(str(path), emu(x), emu(y), height=emu(h))


def movie(slide, video, poster, x, y, w, h):
    """Embed a video sized to the poster's aspect, using the poster as its frame."""
    bx, by, bw, bh = _fit(poster, x, y, w, h)
    return slide.shapes.add_movie(
        str(video), emu(bx), emu(by), emu(bw), emu(bh),
        poster_frame_image=str(poster), mime_type="video/mp4",
    )


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ---------------------------------------------------------------- chrome
def header(slide, title, eyebrow=None, eyebrow_color=ACCENT):
    if eyebrow:
        tf = textbox(slide, ML, 0.34, CW, 0.26)
        para(tf, eyebrow.upper(), size=11.5, color=eyebrow_color, bold=True, first=True)
    tf = textbox(slide, ML, 0.62, CW, 0.52)
    para(tf, title, size=27, color=NAVY, bold=True, first=True, line=1.0)
    rect(slide, ML, 1.235, 1.42, 0.05, fill=ACCENT)


def footer(slide, n, note="SURF 2026  ·  Baaqer Farhat  ·  Caltech ARCL"):
    tf = textbox(slide, ML, FOOTER_Y, CW * 0.75, 0.24)
    para(tf, note, size=9.5, color=GRAY, first=True)
    tf = textbox(slide, SW - MR - 1.0, FOOTER_Y, 1.0, 0.24, align=PP_ALIGN.RIGHT)
    para(tf, str(n), size=9.5, color=GRAY, first=True, align=PP_ALIGN.RIGHT)


def kicker(slide, y, text, color=NAVY, fill=LIGHT, edge=None, size=13.5, h=0.52,
           x=ML, w=CW, bold_prefix=None):
    """Full width emphasis band."""
    rrect(slide, x, y, w, h, radius=0.10, fill=fill,
          line=edge if edge else None, lw=1.3)
    tf = textbox(slide, x + 0.22, y, w - 0.44, h, anchor=MSO_ANCHOR.MIDDLE)
    if bold_prefix:
        rich(tf, [(bold_prefix, {"bold": True, "color": color, "size": size}),
                  (text, {"color": color, "size": size})], first=True)
    else:
        para(tf, text, size=size, color=color, first=True)


def bullets(slide, x, y, w, items, size=14.5, gap=9, color=NAVY, dot=ACCENT,
            line=1.26):
    """items: str, or (bold_lead, rest)."""
    tf = textbox(slide, x, y, w, 0.4)
    for i, it in enumerate(items):
        lead, rest = (it, "") if isinstance(it, str) else it
        chunks = [("\u2022   ", {"color": dot, "bold": True, "size": size}),
                  (lead, {"bold": bool(rest), "color": color, "size": size})]
        if rest:
            chunks.append((rest, {"color": color, "size": size}))
        rich(tf, chunks, first=(i == 0), space_after=gap, line=line)
    return tf


# ---------------------------------------------------------------- video slot
def video_slot(slide, x, y, w, poster, tag, title, caption, source,
               accent=ACCENT, media_h=None, video=None):
    """Labelled video slot.

    Shows a real frame from the rollout it stands for. With SURF_EMBED_VIDEOS=1
    the matching .mp4 is embedded in place of the still.
    """
    bar_h = 0.34
    iw, ih = Image.open(poster).size
    mh = media_h if media_h is not None else w * ih / iw

    # header bar
    rect(slide, x, y, w, bar_h, fill=accent)
    tf = textbox(slide, x + 0.14, y, w - 0.28, bar_h, anchor=MSO_ANCHOR.MIDDLE)
    rich(tf, [(tag + "   ", {"bold": True, "color": WHITE, "size": 10.5}),
              (title, {"color": WHITE, "size": 10.5})], first=True)

    # media well
    rect(slide, x, y + bar_h, w, mh, fill=INK, line=accent, lw=1.25)
    embedded = EMBED and video is not None and Path(video).exists()
    if embedded:
        movie(slide, video, poster, x, y + bar_h, w, mh)
    else:
        picture(slide, poster, x, y + bar_h, w=w, h=mh)
    rect(slide, x, y + bar_h, w, mh, fill=None, line=accent, lw=1.25)

    if not embedded:
        d = 0.44
        cx, cy = x + w / 2 - d / 2, y + bar_h + mh / 2 - d / 2
        ring = shape(slide, MSO_SHAPE.OVAL, cx, cy, d, d, fill=WHITE,
                     line=accent, lw=1.5)
        ring.fill.fore_color.rgb = WHITE
        tri = shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE,
                    cx + d * 0.30, cy + d * 0.25, d * 0.36, d * 0.50, fill=accent)
        tri.rotation = 90

    # caption + source
    cy2 = y + bar_h + mh + 0.10
    tf = textbox(slide, x, cy2, w, 0.62)
    para(tf, caption, size=11.5, color=NAVY, first=True, line=1.24)
    para(tf, source, size=8.5, color=GRAY, font=MONO, space_before=4, line=1.1)
    return cy2 + 0.62


def video_row(slide, y, specs, accent=ACCENT, gap=0.30, w=None, x0=None,
              media_h=None):
    n = len(specs)
    w = w if w else (CW - gap * (n - 1)) / n
    x0 = x0 if x0 is not None else (SW - (w * n + gap * (n - 1))) / 2
    end = y
    for i, s in enumerate(specs):
        end = video_slot(slide, x0 + i * (w + gap), y, w,
                         accent=s.get("accent", accent), media_h=media_h, **{
                             k: v for k, v in s.items() if k != "accent"})
    return end


# ================================================================== slides
def s_title(prs, n):
    s = add_slide(prs)
    rect(s, 0, 0, SW, 2.06, fill=NAVY)
    rect(s, 0, 2.06, SW, 0.055, fill=ACCENT)

    tf = textbox(s, ML, 0.60, CW - 1.2, 1.30)
    para(tf, "Resilient Humanoid Loco-Manipulation", size=34, color=WHITE,
         bold=True, first=True, line=1.06)
    para(tf, "Learned whole-body policies and online layer adaptation "
             "under actuator faults", size=17,
         color=RGBColor(0xA9, 0xC4, 0xDE), space_before=9, line=1.2)

    tf = textbox(s, ML, 2.72, CW, 0.5)
    para(tf, "Baaqer Farhat", size=25, color=NAVY, bold=True, first=True)

    tf = textbox(s, ML, 3.34, 6.2, 1.5)
    para(tf, "Mentor:  Mahdi Taheri", size=14.5, color=NAVY, first=True,
         space_after=5)
    para(tf, "Aerospace Robotics and Control Lab (ARCL)", size=14.5, color=GRAY,
         space_after=5)
    para(tf, "California Institute of Technology", size=14.5, color=GRAY,
         space_after=5)
    para(tf, "Summer Undergraduate Research Fellowship, 2026", size=14.5, color=GRAY)

    # capability strip
    y = 5.24
    items = [
        ("8.7 s", "whole-body box pickup"),
        ("57", "logged hardware runs"),
        ("+53%", "survival under a knee fault"),
        ("50 Hz", "online adaptation, numpy"),
    ]
    w = (CW - 0.30 * 3) / 4
    for i, (big, small) in enumerate(items):
        x = ML + i * (w + 0.30)
        rrect(s, x, y, w, 1.02, radius=0.10, fill=LIGHT)
        rect(s, x, y, 0.055, 1.02, fill=ACCENT)
        tf = textbox(s, x + 0.24, y + 0.15, w - 0.4, 0.42)
        para(tf, big, size=22, color=NAVY, bold=True, first=True, line=1.0)
        tf = textbox(s, x + 0.24, y + 0.60, w - 0.4, 0.36)
        para(tf, small, size=10.5, color=GRAY, first=True, line=1.15)

    notes(s, """
Good afternoon. My name is Baaqer Farhat, and this is work done in ARCL with Mahdi Taheri on
the AgiBot X2 humanoid. The summer covers two whole-body behaviors, taken from motion capture
through to the real robot, and then a controlled actuator fault to test whether the
controller can adapt to damage in the robot itself.
""")
    return s


def s_problem(prs, n):
    s = add_slide(prs)
    header(s, "A policy is only correct for the robot it was trained on",
           "The problem")

    left_w = 6.45
    bullets(s, ML, BODY_TOP + 0.24, left_w, [
        ("Every policy is optimized against one fixed set of dynamics: ",
         "nominal PD gains, a known payload, a particular contact model."),
        ("Real robots leave that set. ",
         "Actuators weaken, payloads change, and the sensors the policy trained "
         "on are not the sensors the hardware has."),
        ("A frozen policy cannot notice. ",
         "It keeps issuing the commands that were optimal for a robot it is no "
         "longer controlling."),
    ], size=15.5, gap=15)

    tf = textbox(s, ML, 4.66, left_w, 1.4)
    para(tf, "Objective", size=12, color=ACCENT, bold=True, first=True,
         space_after=7)
    para(tf, "Behaviors that survive conditions outside the training "
             "distribution, and a controller that responds to changes in the "
             "robot's own dynamics without being told what changed.",
         size=15, color=NAVY, line=1.3)

    # right: what actually broke, from this project's logs
    x = ML + left_w + 0.52
    w = CW - left_w - 0.52
    rrect(s, x, BODY_TOP, w, 4.10, radius=0.055, fill=LIGHT)
    tf = textbox(s, x + 0.30, BODY_TOP + 0.26, w - 0.60, 0.3)
    para(tf, "OBSERVED ON THIS ROBOT", size=11, color=GRAY, bold=True, first=True)

    rows = [
        ("Feet on their edges", "Sim reported contact force, so the reward "
                                "counted the foot as planted.", RED),
        ("Saturated ankle roll", "Frontal-plane correction ran through a joint "
                                 "with 15\u00b0 of travel.", RED),
        ("Wrong base frame", "Torso IMU fed into an observation holosoma builds "
                             "from the pelvis.", RED),
        ("Weakened actuator", "Right knee at 30% stiffness. Survival falls from "
                              "14.7 s to 7.4 s.", AMBER),
    ]
    yy = BODY_TOP + 0.62
    for i, (t, d, c) in enumerate(rows):
        rect(s, x + 0.30, yy + 0.045, 0.045, 0.72, fill=c)
        tf = textbox(s, x + 0.50, yy, w - 0.82, 0.80)
        para(tf, t, size=13, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, d, size=11, color=GRAY, line=1.22)
        yy += 0.87

    footer(s, n)
    notes(s, """
Reinforcement learning gives you a humanoid behavior that works, but you optimize it against
one particular robot: specific actuator gains, a specific payload, a specific contact model.
Real hardware leaves that set almost immediately.

On the right are four failures from this project's run logs. The robot stood on the edges of
its feet. It tried to balance sideways through a joint with fifteen degrees of travel. The
deployment fed the wrong IMU into an observation. And the one introduced deliberately, a
weakened knee.

In every case the policy has no way to know. So the objective is two things: behaviors that
hold up outside the training distribution, and a controller that responds when its own
dynamics change.
""")
    return s


def s_system(prs, n):
    s = add_slide(prs)
    header(s, "From human motion capture to a policy on the robot", "System and workflow")

    y, h = 2.46, 0.94
    w = 1.72
    gap = (CW - 6 * w) / 5
    stages = [
        ("Human MoCap", "OmniRetarget\nCMU", LIGHT, NAVY),
        ("Retargeting", "interaction mesh\n\u2192 31 DOF", LIGHT, NAVY),
        ("Reference clip", "8.7 s pickup\n19 s crawl", LIGHT, NAVY),
        ("RL tracking", "IsaacLab +\nholosoma, PPO", PALE_B, ACCENT),
        ("Export", "self-contained\nnumpy .npz", LIGHT, NAVY),
        ("AgiBot X2", "ROS 2\n50 Hz", PALE_B, ACCENT),
    ]
    xs = [ML + i * (w + gap) for i in range(6)]
    for x, (t, sub, fc, ec) in zip(xs, stages):
        rrect(s, x, y, w, h, radius=0.10, fill=fc, line=ec, lw=1.3)
        tf = textbox(s, x + 0.08, y + 0.13, w - 0.16, 0.28, align=PP_ALIGN.CENTER)
        para(tf, t, size=12.5, color=NAVY, bold=True, first=True,
             align=PP_ALIGN.CENTER, line=1.0)
        tf = textbox(s, x + 0.08, y + 0.44, w - 0.16, 0.42, align=PP_ALIGN.CENTER)
        para(tf, sub, size=10, color=GRAY, first=True, align=PP_ALIGN.CENTER,
             line=1.2)
    for i in range(5):
        arrow(s, xs[i] + w, y + h / 2, xs[i + 1], y + h / 2, color=RULE, lw=1.7)

    # VLM branch, above the robot, deliberately off the control path
    vy = 1.46
    rrect(s, xs[5] - 0.66, vy, w + 0.66, 0.68, radius=0.10, fill=PALE_A,
          line=AMBER, lw=1.4)
    tf = textbox(s, xs[5] - 0.54, vy + 0.09, w + 0.42, 0.52, align=PP_ALIGN.CENTER)
    para(tf, "Head RGB \u2192 VLM", size=11, color=AMBER, bold=True, first=True,
         align=PP_ALIGN.CENTER, line=1.0)
    para(tf, "scene description, open loop", size=9, color=AMBER, space_before=3,
         align=PP_ALIGN.CENTER)
    arrow(s, xs[5] + w / 2, y, xs[5] + w / 2, vy + 0.68, color=AMBER, lw=1.5,
          dashed=True)

    # adaptation branch, between the policy and the actuators
    ay = 3.86
    rrect(s, xs[4] - 0.22, ay, w + 0.44, 0.72, radius=0.10, fill=PALE_G,
          line=GREEN, lw=1.5)
    tf = textbox(s, xs[4] - 0.14, ay + 0.10, w + 0.28, 0.56, align=PP_ALIGN.CENTER)
    para(tf, "Online layer adaptation", size=11, color=GREEN, bold=True,
         first=True, align=PP_ALIGN.CENTER, line=1.0)
    para(tf, "50 Hz  ·  numpy  ·  5.25 ms / step", size=9, color=GREEN,
         align=PP_ALIGN.CENTER, space_before=3)
    arrow(s, xs[4] + w / 2, y + h, xs[4] + w / 2, ay, color=GREEN, lw=1.6)
    arrow(s, xs[4] + w + 0.22, ay + 0.36, xs[5] + w / 2, ay + 0.36, color=GREEN, lw=1.6)
    arrow(s, xs[5] + w / 2, ay + 0.36, xs[5] + w / 2, y + h, color=GREEN, lw=1.6)

    # hardware feedback loop
    fy = 5.16
    fx = xs[5] + w - 0.22
    arrow(s, fx, y + h, fx, fy, color=RED, lw=1.5, dashed=True)
    arrow(s, fx, fy, xs[3] + w / 2, fy, color=RED, lw=1.5, dashed=True)
    arrow(s, xs[3] + w / 2, fy, xs[3] + w / 2, y + h, color=RED, lw=1.5, dashed=True)
    tf = textbox(s, xs[3], fy + 0.10, fx - xs[3], 0.3, align=PP_ALIGN.CENTER)
    para(tf, "57 logged hardware runs \u2192 new reward terms and observation fixes",
         size=11, color=RED, italic=True, first=True, align=PP_ALIGN.CENTER)

    # training detail band
    rrect(s, ML, 5.86, CW, 0.50, radius=0.10, fill=LIGHT)
    tf = textbox(s, ML + 0.24, 5.86, CW - 0.48, 0.50, anchor=MSO_ANCHOR.MIDDLE)
    rich(tf, [
        ("4096 ", {"bold": True}), ("parallel envs      ", {"color": GRAY}),
        ("200 Hz ", {"bold": True}), ("physics / ", {"color": GRAY}),
        ("50 Hz ", {"bold": True}), ("control      ", {"color": GRAY}),
        ("blind actor MLP 512-256-128 ", {"bold": True}),
        ("(ELU)      ", {"color": GRAY}),
        ("DR: ", {"bold": True}),
        ("box mass, friction, PD gain, action delay, encoder noise, pushes",
         {"color": GRAY}),
    ], size=11, first=True)

    footer(s, n)
    notes(s, """
The pipeline is the same for both behaviors. Human motion capture, retargeted onto the thirty
one degree of freedom X2, gives a reference clip. A policy learns to track it with PPO in
IsaacLab, four thousand environments in parallel. The actor is blind: no cameras, only
proprioception and a clock into the reference. It exports to numpy and runs on the robot over
ROS 2 at fifty hertz.

Two things to notice. The red loop is real: fifty seven hardware runs fed back into the
reward design. In green is where the adaptation sits, between the policy and the actuators.
The VLM in orange is off to the side on purpose, because it is not closed on control.
""")
    return s


def s_box_task(prs, n):
    s = add_slide(prs)
    header(s, "Box pickup: one policy for the whole 8.7 second motion",
           "Behavior 1  ·  loco-manipulation")

    tf = textbox(s, ML, 1.44, 8.5, 0.62)
    para(tf, "From the default standing pose the robot bends to a 45 cm box, "
             "squeezes it between two rigid fingerless palms, lifts it to chest "
             "height, carries it 1.4 m, and sets it down over planted feet.",
         size=13.5, color=NAVY, first=True, line=1.26)

    x = ML + 8.86
    w = CW - 8.86
    rrect(s, x, 1.38, w, 0.78, radius=0.10, fill=PALE_R, line=RED, lw=1.3)
    tf = textbox(s, x + 0.18, 1.38, w - 0.36, 0.78, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "No fingers, no cameras, and the payload moves the "
             "centre of mass while the robot is bent over.",
         size=11, color=RED, first=True, line=1.22)

    video_row(s, 2.20, [
        dict(poster=POSTERS / "box1_v8.png", tag="VIDEO 1",
             title="Early  ·  v8, iter 8000",
             caption="Balances through the full motion but pantomimes the lift. "
                     "The box never leaves the floor.",
             source="x2_box_v8_iter8000_full.mp4",
             video=VIDEOS / "x2_box_v8_iter8000_full.mp4"),
        dict(poster=POSTERS / "box2_v10.png", tag="VIDEO 2",
             title="Intermediate  ·  v10, iter 28000",
             caption="First complete pickup, carry and set-down in simulation "
                     "after the two-hand grasp reward.",
             source="x2_box_v10_iter28000_final.mp4",
             video=VIDEOS / "x2_box_v10_iter28000_final.mp4"),
        dict(poster=POSTERS / "box3_v33.png", tag="VIDEO 3",
             title="Current  ·  v33, iter 253000",
             caption="Flat planted feet, quiet waist, returns fully upright. "
                     "This is the deployed policy.",
             source="x2_box_v33_waist_track_iter253000.mp4",
             video=VIDEOS / "x2_box_v33_waist_track_iter253000.mp4",
             accent=GREEN),
    ])

    # phase strip of the reference motion
    py, ph = 5.90, 0.46
    phases = [
        ("hold upright", 1.0, LIGHT, NAVY),
        ("bend down", 1.5, PALE_B, NAVY),
        ("grasp + lift", 0.5, ACCENT, WHITE),
        ("carry 1.4 m", 2.8, PALE_B, NAVY),
        ("set down", 2.2, PALE_B, NAVY),
        ("stand up", 0.7, LIGHT, NAVY),
    ]
    total = sum(p[1] for p in phases)
    x = ML
    for label, dur, fc, tc in phases:
        pw = CW * dur / total
        rect(s, x, py, pw, ph, fill=fc, line=WHITE, lw=1.5)
        tf = textbox(s, x + 0.04, py, pw - 0.08, ph, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.CENTER)
        para(tf, label, size=10, color=tc, bold=True, first=True,
             align=PP_ALIGN.CENTER)
        x += pw
    tf = textbox(s, ML, py + ph + 0.09, CW, 0.28)
    para(tf, "0 s", size=9.5, color=GRAY, first=True)
    tf = textbox(s, ML, py + ph + 0.09, CW, 0.28, align=PP_ALIGN.RIGHT)
    para(tf, "8.7 s", size=9.5, color=GRAY, first=True, align=PP_ALIGN.RIGHT)
    tf = textbox(s, ML, py - 0.32, CW, 0.28)
    para(tf, "One continuous reference motion, no state machine",
         size=10.5, color=GRAY, first=True)

    footer(s, n)
    notes(s, """
The first behavior is a box pickup, one policy for the whole motion. No state machine. Bend,
grasp, lift, carry, set down, stand up, in eight point seven seconds.

Three things make it hard. The X2 has no fingers, so the grasp is a squeeze between two rigid
palms. The policy is blind. And once the box comes up the centre of mass moves forward while
the robot is already bent over, which is when it has the least margin.

Left to right. Version eight balances through the whole motion, but watch the box. It never
leaves the ground: the policy found it could collect most of the tracking reward by miming
the motion. The middle, version ten, once a grasp reward was added, is the first time it
completes the task. On the right is the policy running on hardware today.
""")
    return s


def s_box_lessons(prs, n):
    s = add_slide(prs)
    header(s, "Most of the work was finding out why it did not work",
           "Behavior 1  ·  development")

    rows = [
        ("v8", "Forgot the walk-up",
         "Sampler mixed raw failure counts, so the uniform term was a no-op.",
         "Mix normalized probabilities."),
        ("v10", "Pantomimed the lift",
         "A mean over hands let one palm mask the other hovering away.",
         "Grasp reward as a product over both hands."),
        ("v13", "Frozen hip, impossible waist",
         "Off-by-one in the retargeter clamped every joint to the previous "
         "joint's limits.",
         "Index limits by joint type, retrain."),
        ("v24", "Crab-walked the carry",
         "The demonstrator turned 65 degrees and shuffled his feet 30 cm.",
         "Straighten the carry, pin the feet."),
    ]
    y = BODY_TOP - 0.04
    rh = 0.62
    for i, (tag, sym, why, fix) in enumerate(rows):
        if i % 2 == 0:
            rect(s, ML, y - 0.03, CW, rh, fill=RGBColor(0xF7, 0xFA, 0xFC))
        tf = textbox(s, ML + 0.06, y + 0.11, 0.62, 0.3)
        para(tf, tag, size=13, color=ACCENT, bold=True, first=True)
        tf = textbox(s, ML + 0.72, y + 0.09, 2.62, 0.44)
        para(tf, sym, size=12.5, color=RED, bold=True, first=True, line=1.14)
        tf = textbox(s, ML + 3.46, y + 0.09, 5.16, 0.50)
        para(tf, why, size=11, color=GRAY, first=True, line=1.18)
        tf = textbox(s, ML + 8.78, y + 0.09, CW - 8.84, 0.50)
        para(tf, fix, size=11, color=GREEN, first=True, line=1.18)
        y += rh

    tf = textbox(s, ML, y + 0.08, CW, 0.30)
    rich(tf, [("Then hardware found what simulation could not.  ",
               {"bold": True, "color": NAVY, "size": 12.5}),
              ("A foot on its edge still reports contact force, and a tilted "
               "stationary foot has zero slip.",
               {"color": GRAY, "size": 12.5})], first=True, line=1.2)

    picture(s, ASSETS / "fig_foot.png", ML, y + 0.46, w=CW)

    footer(s, n)
    notes(s, """
This is where the actual work was. Each row is a failure that had to be diagnosed before the
next version could improve on it. Version eight kept forgetting the walk-up because the
sampler mixed raw failure counts, which scale with the environment count, so the uniform term
was effectively zero. Ten is the pantomime you just saw.

The plots are what simulation could not find on its own. On the left, foot tilt across five
checkpoints: the right foot sits past twenty degrees, standing on its edge, and no reward
objected, because the engine still reports contact force on an edge.

On the right is the loop that followed. The feet skated, so a slip penalty went in and foot
travel dropped from twenty eight centimetres to half a centimetre. Then they stepped instead
of sliding, which is free under a slip penalty, so a contact loss term went in. Then the
waist rewards out-competed the slip penalty and it had to double again.
""")
    return s


def s_box_diagnosis(prs, n):
    s = add_slide(prs)
    header(s, "The hardware logs pointed at a single joint",
           "Behavior 1  ·  sim to real")

    tf = textbox(s, ML, 1.40, CW, 0.34)
    para(tf, "Version 31 was the best policy in simulation and it fell on the "
             "robot, 13 of 15 runs, always at the same point in the motion.",
         size=13.5, color=NAVY, first=True, line=1.22)

    picture(s, ASSETS / "fig_joint_error.png", ML, 1.84, w=CW)

    y = 5.24
    cards = [
        ("Symptom", RED, PALE_R,
         "Torso pitched forward and collapsed at t \u2248 8.2 s, entering the "
         "carry. Leg tracking error stayed near 5\u00b0 throughout."),
        ("Cause", AMBER, PALE_A,
         "The actor drove the waist pitch command to \u2212131\u00b0 while the "
         "reference clip asks for at most +18\u00b0. Nothing in the reward "
         "priced the waist."),
        ("Fix, v33", GREEN, PALE_G,
         "Waist pitch and yaw tracking terms, exponential plus L2. Command "
         "back inside the clip and hardware waist error 13.1\u00b0 to 3.3\u00b0."),
    ]
    w = (CW - 2 * 0.28) / 3
    for i, (t, c, fc, body) in enumerate(cards):
        x = ML + i * (w + 0.28)
        rrect(s, x, y, w, 1.42, radius=0.06, fill=fc, line=c, lw=1.3)
        tf = textbox(s, x + 0.22, y + 0.14, w - 0.44, 0.26)
        para(tf, t.upper(), size=10.5, color=c, bold=True, first=True)
        tf = textbox(s, x + 0.22, y + 0.46, w - 0.44, 0.86)
        para(tf, body, size=11.5, color=NAVY, first=True, line=1.24)

    footer(s, n)
    notes(s, """
This is the sim to real gap, and the clearest example of the whole workflow. Version thirty
one was the best checkpoint in simulation. On the robot it fell thirteen times out of
fifteen, always at about eight seconds, always the same way: the torso pitched forward and
did not come back.

The left plot ranks every joint by tracking error across those runs. The legs sit around five
degrees, which is why the first guess, a leg problem, was wrong. One bar is at eighty four
degrees, and it is the waist pitch.

The middle plot shows what that means. The dashed line is what the reference clip asks for,
at most eighteen degrees forward. The red line is what the policy commanded: minus one
hundred and thirty one degrees, far outside anything in the clip, so the joint saturates.
Nothing in the reward priced the waist, so the policy used it as a dumping ground and
simulation absorbed it.

The right plot is the same measurement after adding waist tracking terms. The command stays
inside the clip and hardware waist error drops from thirteen degrees to three.
""")
    return s


def s_crawl(prs, n):
    s = add_slide(prs)
    header(s, "Crawling: a prone gait needs its own reward set",
           "Behavior 2  ·  locomotion")

    tf = textbox(s, ML, 1.44, 7.35, 0.62)
    para(tf, "A whole-body prone crawl up a slope, retargeted from a CMU motion "
             "capture clip and trained in the same holosoma stack. The robot moves "
             "on palms, knees and feet.",
         size=13.5, color=NAVY, first=True, line=1.26)

    x = ML + 7.66
    w = CW - 7.66
    rrect(s, x, 1.36, w, 1.00, radius=0.10, fill=PALE_A, line=AMBER, lw=1.3)
    tf = textbox(s, x + 0.18, 1.36, w - 0.36, 1.00, anchor=MSO_ANCHOR.MIDDLE)
    rich(tf, [("The upright reward suite is hostile to crawling. ",
               {"bold": True, "color": AMBER, "size": 10.5}),
              ("Foot slip penalizes the limb sliding a crawl needs, and contact "
               "loss demands both feet stay loaded. The task needs its own "
               "reward, contact and observation set.",
               {"color": AMBER, "size": 10.5})], first=True, line=1.2)

    video_row(s, 2.40, [
        dict(poster=POSTERS / "crawl1_v2.png", tag="VIDEO 1",
             title="Early  ·  v2, iter 14500",
             caption="Crawls for the first time. Tracking is loose and the robot "
                     "drifts off the slope.",
             source="x2_crawl_slope_v2_palmflat_iter14500.mp4",
             video=VIDEOS / "x2_crawl_slope_v2_palmflat_iter14500.mp4"),
        dict(poster=POSTERS / "crawl2_v3.png", tag="VIDEO 2",
             title="Intermediate  ·  v3, iter 49999",
             caption="Retuned tracking weights. This checkpoint was exported and "
                     "run on the physical X2.",
             source="x2_crawl_slope_v3_tracking_iter49999.mp4",
             video=VIDEOS / "x2_crawl_slope_v3_tracking_iter49999.mp4"),
        dict(poster=POSTERS / "crawl3_v5.png", tag="VIDEO 3",
             title="Current  ·  v5, iter 86000",
             caption="Position weight raised to 3.0. Mean episode 7.6 s, the best "
                     "climb so far in simulation.",
             source="x2_crawl_slope_v5_track_xyz_iter86000.mp4",
             video=VIDEOS / "x2_crawl_slope_v5_track_xyz_iter86000.mp4",
             accent=GREEN),
    ])

    kicker(s, 6.06,
           "  v4 over-smoothed the gait until it could not generate thrust, and v6 "
           "multi-clip stand-up training regressed the climb from 24.9 to 16.2 mean "
           "reward. Best episode is still 7.6 s of a 19 s clip.",
           bold_prefix="Honest limitations.", fill=WHITE, edge=RULE, size=12, h=0.60)

    footer(s, n)
    notes(s, """
The second behavior is crawling. Same stack, but prone, climbing a slope on palms and knees.

The technical point is the orange box. The first attempt loaded the slope into the existing
pickup task and failed completely, because the upright reward suite fights a crawl. There is
a foot slip penalty, and sliding your limbs along the ground is exactly what crawling is. So
the task got its own reward set, knees, elbows and palms added to the allowed contacts, and
projected gravity in the actor so the policy can tell it is belly down.

The middle video is the checkpoint that ran on the physical robot. On the limits: version
four over smoothed the gait until it could not climb, version six made the crawl worse, and
the best policy sustains seven and a half seconds of a nineteen second clip. Real, but less
mature than the pickup.
""")
    return s


def s_why_adapt(prs, n):
    s = add_slide(prs)
    header(s, "A good policy is still a fixed policy",
           "Transition  ·  why adaptation")

    tf = textbox(s, ML, 1.44, CW, 0.4)
    para(tf, "To test that, the right knee actuator is scaled to 30% of nominal "
             "PD stiffness. It still moves, but it sags under load.",
         size=15, color=NAVY, first=True)

    y, h = 2.28, 2.20
    w = (CW - 3 * 0.42) / 4
    xs = [ML + i * (w + 0.42) for i in range(4)]
    panels = [
        ("Nominal robot", "The frozen policy completes the pickup on all six seeds.",
         "14.68 s", NAVY, WHITE, RULE),
        ("Actuator fault", "Right knee drops to 30% of nominal PD stiffness.",
         "not announced", AMBER, PALE_A, AMBER),
        ("Frozen policy", "Sags under the payload, loses the set-down, falls early.",
         "7.43 s", RED, PALE_R, RED),
        ("Adapted policy", "Recovers the set-down on all six seeds.",
         "11.36 s", GREEN, PALE_G, GREEN),
    ]
    for x, (t, body, tag, c, fc, ec) in zip(xs, panels):
        rrect(s, x, y, w, h, radius=0.07, fill=fc, line=ec, lw=1.6)
        tf = textbox(s, x + 0.18, y + 0.24, w - 0.36, 0.34, align=PP_ALIGN.CENTER)
        para(tf, t, size=15, color=c, bold=True, first=True, align=PP_ALIGN.CENTER)
        tf = textbox(s, x + 0.20, y + 0.78, w - 0.40, 0.72, align=PP_ALIGN.CENTER)
        para(tf, body, size=11.5, color=NAVY, first=True, align=PP_ALIGN.CENTER,
             line=1.28)
        tf = textbox(s, x + 0.14, y + h - 0.62, w - 0.28, 0.42, align=PP_ALIGN.CENTER)
        para(tf, tag, size=19 if tag.endswith("s") else 13, color=c, bold=True,
             first=True, align=PP_ALIGN.CENTER)
    for i in range(3):
        arrow(s, xs[i] + w + 0.05, y + h / 2, xs[i + 1] - 0.05, y + h / 2,
              color=NAVY, lw=2.0)

    kicker(s, 4.92,
           "  Neither the policy nor the adapter is told that a fault occurred. "
           "Both keep using the nominal gains. There is no detector and no trigger.",
           bold_prefix="The fault is never announced.", fill=LIGHT, size=13.5, h=0.56)

    tf = textbox(s, ML, 5.76, CW, 0.9)
    para(tf, "This is the regime the method is actually for. The same adaptation "
             "applied to a healthy robot spends balance margin the task needed and "
             "makes things worse. That measurement is in backup.",
         size=12.5, color=GRAY, first=True, italic=True, line=1.3)

    footer(s, n)
    notes(s, """
Both behaviors work, but they are fixed policies, and that brings us to the main question.

The fault has to be realistic, so the right knee stiffness is scaled to thirty percent of
nominal. It still moves, but it sags whenever it is loaded, which for a robot bent over
holding a box is most of the motion.

Healthy, the frozen policy is essentially perfect. Add the fault and it drops to seven point
four and becomes erratic. With adaptation it comes back to eleven point four.

The point to land is the grey band. Nothing tells the policy a fault happened. No detector,
no trigger. The adapter runs from the first control step and only ever sees tracking error.
""")
    return s


def s_method(prs, n):
    s = add_slide(prs)
    header(s, "Adapt one layer, online, from tracking error alone",
           "Adaptation  ·  method")

    tf = textbox(s, ML, 1.42, CW, 0.34)
    rich(tf, [("Ported from ", {"color": GRAY, "size": 12.5}),
              ("Taheri, Chung and Hadaegh, ACC 2026, ", {"color": NAVY, "size": 12.5}),
              ("\u201cClosing the Loop Inside Neural Networks: Causality-Guided Layer "
               "Adaptation for Fault Recovery Control\u201d", {"italic": True,
                                                               "color": GRAY,
                                                               "size": 12.5}),
              ("  \u2014  moved from its MuJoCo testbed into our IsaacLab environment "
               "on the v31 policy.", {"color": GRAY, "size": 12.5})],
         first=True, line=1.2)

    # --- network diagram
    y, h = 2.44, 0.86
    w = 1.86
    gap = 0.40
    x0 = ML + 0.30
    layers = [
        ("W\u2080", "164 \u2192 512", False),
        ("W\u2081", "512 \u2192 256", False),
        ("W\u2082", "256 \u2192 128", True),
        ("W\u2083", "128 \u2192 31", False),
    ]
    xs = [x0 + i * (w + gap) for i in range(4)]
    for x, (lab, sub, adapt) in zip(xs, layers):
        rrect(s, x, y, w, h, radius=0.10,
              fill=PALE_G if adapt else LIGHT,
              line=GREEN if adapt else RULE, lw=2.2 if adapt else 1.2)
        tf = textbox(s, x + 0.08, y + 0.12, w - 0.16, 0.30, align=PP_ALIGN.CENTER)
        para(tf, lab, size=14, color=GREEN if adapt else NAVY, bold=True,
             first=True, align=PP_ALIGN.CENTER, line=1.0)
        tf = textbox(s, x + 0.08, y + 0.46, w - 0.16, 0.28, align=PP_ALIGN.CENTER)
        para(tf, sub, size=10.5, color=GREEN if adapt else GRAY, first=True,
             align=PP_ALIGN.CENTER)
    for i in range(3):
        arrow(s, xs[i] + w, y + h / 2, xs[i + 1], y + h / 2, color=RULE, lw=1.5)

    tf = textbox(s, ML - 0.30, y + 0.18, 0.58, 0.5, align=PP_ALIGN.CENTER)
    para(tf, "obs\n164", size=10, color=GRAY, first=True, align=PP_ALIGN.CENTER,
         line=1.15)
    arrow(s, ML + 0.06, y + h / 2, x0, y + h / 2, color=RULE, lw=1.5)
    arrow(s, xs[3] + w, y + h / 2, xs[3] + w + 0.44, y + h / 2, color=RULE, lw=1.5)
    tf = textbox(s, xs[3] + w + 0.10, y + 0.18, 1.1, 0.5)
    para(tf, "31 joint\ntargets", size=10, color=GRAY, first=True, line=1.15)

    # frozen / adapted captions
    for i, x in enumerate(xs):
        if i == 2:
            continue
        tf = textbox(s, x, y - 0.30, w, 0.26, align=PP_ALIGN.CENTER)
        para(tf, "FROZEN", size=9.5, color=GRAY, bold=True, first=True,
             align=PP_ALIGN.CENTER)
    tf = textbox(s, xs[2] - 0.92, y - 0.30, w + 1.84, 0.26, align=PP_ALIGN.CENTER)
    para(tf, "ADAPTED  ·  32,768 of 251,776 weights", size=10, color=GREEN,
         bold=True, first=True, align=PP_ALIGN.CENTER)

    # feedback path
    ey = 3.80
    rrect(s, xs[1], ey, w * 3 + gap * 2, 0.54, radius=0.10, fill=PALE_G,
          line=GREEN, lw=1.5)
    tf = textbox(s, xs[1] + 0.16, ey, w * 3 + gap * 2 - 0.32, 0.54,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    rich(tf, [("e = q \u2212 q", {"bold": True, "color": GREEN, "size": 14}),
              ("ref", {"bold": True, "color": GREEN, "size": 9}),
              ("     joint tracking error, measured every 20 ms",
               {"color": GREEN, "size": 11.5})],
         first=True, align=PP_ALIGN.CENTER)
    arrow(s, xs[3] + w * 0.55, y + h, xs[3] + w * 0.55, ey + 0.27, color=GREEN, lw=1.6)
    arrow(s, xs[3] + w * 0.55, ey + 0.27, xs[1] + w * 3 + gap * 2, ey + 0.27,
          color=GREEN, lw=1.6)
    arrow(s, xs[2] + w / 2, ey, xs[2] + w / 2, y + h, color=GREEN, lw=1.6)

    eq = ASSETS / "eq_update.png"
    if eq.exists():
        picture(s, eq, ML + 0.30, 4.72, w=3.40)
    tf = textbox(s, ML + 0.30, 5.28, 3.60, 0.3)
    para(tf, "Lyapunov update, leakage toward the trained weights", size=9.5,
         color=GRAY, first=True)

    # right column facts
    cx = ML + 4.44
    cwid = CW - 4.44
    facts = [
        ("Adapted", "one weight matrix, 13% of the network"),
        ("Frozen", "every other layer, the reward, the reference clip"),
        ("Needs", "no plant model, no fault detector, no offline phase"),
        ("Cost", "1.93 ms forward + 3.31 ms update = 26% of the 20 ms budget"),
        ("Mask", "which joints enter e. Waist only, against the paper's "
                 "hips, knees, ankles and waist"),
        ("Safety", "resets to W\u2080 every episode; a norm bound trips a hard revert"),
    ]
    yy = 4.62
    for k, v in facts:
        tf = textbox(s, cx, yy, cwid, 0.32)
        rich(tf, [(k + "   ", {"bold": True, "color": ACCENT, "size": 11}),
                  (v, {"color": NAVY, "size": 11})], first=True, line=1.16)
        yy += 0.305

    kicker(s, 6.30,
           "  the shipped leak pulls the layer toward the origin rather than the "
           "trained weights, eroding 13.7% of its norm per episode even with perfect "
           "tracking. Fixing it does not change the outcome, but it is a real bug.",
           bold_prefix="Defect found and fixed:", fill=PALE_R, edge=RED,
           color=NAVY, size=11.5, h=0.52)

    footer(s, n)
    notes(s, """
The method is from an ACC paper by my mentor, with Soon Jo Chung and Fred Hadaegh. The work
here was taking it out of the MuJoCo testbed it shipped with, getting it running in our
IsaacLab environment on a real trained policy, and finding out whether the result holds.

The policy is four layers. The adapter changes exactly one weight matrix, about thirteen
percent of the parameters. Everything else is frozen. No retraining, no replay, no model of
the plant, and no fault detector.

The only input is the green path: joint tracking error, every twenty milliseconds, pushed
back through the layers above the adapted one to drive a Lyapunov update, with leakage
pulling the weights toward the trained values.

It costs about five milliseconds per step in numpy, which is why the same code runs on the
robot, and it never persists across episodes. The red band is a defect found and fixed along
the way: the leakage pulled weights toward zero rather than toward the trained values.
""")
    return s


def s_results(prs, n):
    s = add_slide(prs)
    header(s, "Under the fault, adaptation recovers the task",
           "Adaptation  ·  results", eyebrow_color=GREEN)

    picture(s, ASSETS / "fig_fault_results.png", ML - 0.10, 1.42, w=6.30, h=4.34)

    x = ML + 6.42
    w = CW - 6.42

    # table
    ty = 1.50
    rows = [
        ("", "Survival", "Set-down", "Lifted"),
        ("Frozen baseline", "7.43 \u00b1 2.76 s", "4 / 6", "5 / 6"),
        ("Waist-only adaptation", "11.36 \u00b1 0.42 s", "6 / 6", "6 / 6"),
    ]
    colx = [0.0, 2.36, 3.86, 4.78]
    colw = [2.36, 1.50, 0.92, 0.89]
    rh = [0.36, 0.56, 0.56]
    yy = ty
    for ri, row in enumerate(rows):
        if ri == 0:
            rect(s, x, yy, w, rh[ri], fill=NAVY)
        elif ri == 2:
            rect(s, x, yy, w, rh[ri], fill=PALE_G)
        else:
            rect(s, x, yy, w, rh[ri], fill=RGBColor(0xF7, 0xFA, 0xFC))
        for ci, cell in enumerate(row):
            al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            tf = textbox(s, x + colx[ci] + (0.16 if ci == 0 else 0), yy,
                         colw[ci], rh[ri], anchor=MSO_ANCHOR.MIDDLE, align=al)
            if ri == 0:
                para(tf, cell, size=10.5, color=WHITE, bold=True, first=True, align=al)
            else:
                c = GREEN if ri == 2 else RED
                para(tf, cell, size=13 if ci else 12.5,
                     color=c if ci else NAVY, bold=(ci > 0), first=True, align=al)
        yy += rh[ri]
    rect(s, x, ty, w, sum(rh), fill=None, line=RULE, lw=1.1)

    tf = textbox(s, x, yy + 0.10, w, 0.3)
    para(tf, "6 seeds  ·  no domain randomization  ·  observation noise off  ·  "
             "v31 policy", size=9.5, color=GRAY, first=True)

    # findings
    findings = [
        ("Complete separation.", " Every adapted seed outlasts every frozen seed. "
                                 "Exact one-sided Mann-Whitney p = 0.0011, the floor "
                                 "for six against six."),
        ("Consistency, not just mean.", " Spread collapses from \u00b12.76 s to "
                                        "\u00b10.42 s. The seed where the frozen "
                                        "policy crumples at 1.3 s is carried to 10.4 s."),
        ("The task comes back.", " Frozen reaches the set-down on 4 of 6. "
                                 "Adaptation reaches it on all 6."),
    ]
    yy = yy + 0.46
    for lead, rest in findings:
        tf = textbox(s, x, yy, w, 0.62)
        rich(tf, [("\u2022  ", {"color": GREEN, "bold": True, "size": 12}),
                  (lead, {"bold": True, "color": NAVY, "size": 12}),
                  (rest, {"color": NAVY, "size": 12})], first=True, line=1.22)
        yy += 0.68

    rrect(s, x, 5.56, w, 1.32, radius=0.06, fill=PALE_A, line=AMBER, lw=1.4)
    tf = textbox(s, x + 0.22, 5.68, w - 0.44, 1.14)
    para(tf, "THE MASK DECIDES, NOT THE GAIN", size=10, color=AMBER, bold=True,
         first=True, space_after=5)
    para(tf, "Waist-only runs the same 3e-4 gain that is catastrophic with the "
             "paper's legs-and-waist mask: 2.63 s and never lifts the box. The legs "
             "are where balance lives on a floating base, so regulating their "
             "tracking error fights the balance the policy is maintaining.",
         size=10, color=NAVY, line=1.22)

    footer(s, n)
    notes(s, """
This is the main result, and it is worth stating precisely. Six seeds, no randomization,
noise off, and the baseline is the same numpy policy with the gain set to zero, so the only
difference between the arms is the adaptation itself.

The dashed line is the healthy robot at fourteen point seven seconds. Red is the frozen
policy under the fault, at seven point four, and look at the spread. One seed dies at one
point three. Green is the same policy with waist adaptation, at eleven point four.

What matters is not the mean, it is the arrow between them. There is no overlap between the
seeds at all: the worst adapted run beats the best frozen run. The spread also collapses from
two point eight seconds to point four, so adaptation removes the catastrophic runs rather than
adding time on average. And the task comes back, from four of six reaching the set down to
all six.

Now the orange box, the unexpected finding. The mask is which joints feed the error signal.
The paper's default includes the legs, and with that mask, at the identical gain, the robot
is on the floor in two and a half seconds having never lifted the box. On a floating base the
legs are where balance lives, so regulating their error fights the balance the policy is
already maintaining.
""")
    return s


def s_adapt_videos(prs, n):
    s = add_slide(prs)
    header(s, "Same policy, same fault, one difference",
           "Adaptation  ·  side by side", eyebrow_color=GREEN)

    tf = textbox(s, ML, 1.42, CW, 0.34)
    para(tf, "All three clips run the v31 box-pickup policy. The only change "
             "between the middle and the right is whether the adaptation update "
             "is switched on.",
         size=13.5, color=NAVY, first=True)

    video_row(s, 1.94, [
        dict(poster=POSTERS / "adapt1_nominal.png", tag="VIDEO 1",
             title="Nominal, frozen",
             caption="Healthy robot. Completes the pickup, carry and set-down.  "
                     "14.68 s, 6/6 seeds.",
             source="x2_box_v31_flatfoot_iter202500.mp4",
             video=VIDEOS / "x2_box_v31_flatfoot_iter202500.mp4",
             accent=NAVY),
        dict(poster=POSTERS / "adapt2_fault_frozen.png", tag="VIDEO 2",
             title="Knee fault, frozen",
             caption="Right knee at 30%. Sags under the payload and falls before "
                     "the set-down.  7.43 s.",
             source="isaac_fault_knee03_frozen_vs_waistadapt.mp4  (left half)",
             video=CLIPS / "fault_knee03_frozen.mp4",
             accent=RED),
        dict(poster=POSTERS / "adapt3_fault_adapt.png", tag="VIDEO 3",
             title="Knee fault, waist adaptation",
             caption="Same fault, adaptation on. Holds the box and reaches the "
                     "set-down.  11.36 s.",
             source="isaac_fault_knee03_frozen_vs_waistadapt.mp4  (right half)",
             video=CLIPS / "fault_knee03_waistadapt.mp4",
             accent=GREEN),
    ], w=3.30, gap=0.40, media_h=3.26)

    footer(s, n)
    notes(s, """
Same result as a picture.

On the left, the healthy robot with the frozen policy. That is what the behavior should look
like. In the middle, the identical policy with the knee at thirty percent. Watch the right
side. It sags as the box loads the weakened knee, the torso pitches too far, and it goes
down before the set down.

On the right, same fault, same starting weights, and the only difference is that the
adaptation update is running. It holds the box and reaches the set down.

Nothing was retrained. The weights on the right started identical to the middle. The
difference accumulated during the ten seconds you are watching.
""")
    return s


def s_vlm(prs, n):
    s = add_slide(prs)
    header(s, "VLM perception: built and benchmarked, not yet closed on control",
           "Perception  ·  longer term autonomy", eyebrow_color=AMBER)

    tf = textbox(s, ML, 1.40, CW, 0.36)
    para(tf, "Adaptation handles faults the robot can feel in its own tracking "
             "error. Deciding what to do next needs the robot to understand the "
             "scene, which is where a vision language model fits.",
         size=12.5, color=NAVY, first=True, line=1.24)

    # pipeline, full width
    y, h = 2.02, 0.74
    w = (CW - 3 * 0.30) / 4
    gap = 0.30
    nodes = [
        ("Orbbec Gemini 335", "head RGB, 1280 \u00d7 720", LIGHT, NAVY),
        ("Camera server on robot", "JPEG over HTTP, port 8099", LIGHT, NAVY),
        ("Ollama on workstation", "Qwen2.5-VL 7B / 3B, Moondream", PALE_A, AMBER),
        ("Scene description", "objects, people, coarse position", PALE_A, AMBER),
    ]
    xs = [ML + i * (w + gap) for i in range(4)]
    for x, (t, sub, fc, ec) in zip(xs, nodes):
        rrect(s, x, y, w, h, radius=0.10, fill=fc, line=ec, lw=1.3)
        tf = textbox(s, x + 0.12, y + 0.12, w - 0.24, 0.28, align=PP_ALIGN.CENTER)
        para(tf, t, size=11.5, color=NAVY, bold=True, first=True,
             align=PP_ALIGN.CENTER, line=1.0)
        tf = textbox(s, x + 0.08, y + 0.42, w - 0.16, 0.28, align=PP_ALIGN.CENTER)
        para(tf, sub, size=9.5, color=GRAY, first=True, align=PP_ALIGN.CENTER)
    for i in range(3):
        arrow(s, xs[i] + w, y + h / 2, xs[i + 1], y + h / 2, color=AMBER, lw=1.5)

    picture(s, ASSETS / "fig_vlm_latency.png", ML - 0.06, 3.02, w=5.45, h=2.42)
    tf = textbox(s, ML + 0.10, 5.56, 5.3, 0.40)
    para(tf, "Warm, end to end, RTX 5070 8 GB. 12 frames from one 15 s clip, "
             "identical prompt, temperature 0.2.",
         size=9.5, color=GRAY, first=True, line=1.2)

    # status columns
    cx = ML + 5.72
    cwid = CW - 5.72
    cols = [
        ("COMPLETED", GREEN, PALE_G, [
            "Head camera streamed off the robot",
            "Three models served locally through Ollama",
            "Latency and description quality measured",
        ]),
        ("CURRENTLY TESTING", AMBER, PALE_A, [
            "Which model is accurate enough at a usable rate",
        ]),
        ("NEXT STEP", ACCENT, PALE_B, [
            "Free text into a structured signal the controller "
            "can consume",
        ]),
    ]
    lines = [3, 1, 2]   # rendered line count per column, incl. wraps
    yy = 3.02
    for (title, c, fc, items), nl in zip(cols, lines):
        hgt = 0.40 + 0.235 * nl + 0.16
        rrect(s, cx, yy, cwid, hgt, radius=0.06, fill=fc, line=c, lw=1.3)
        tf = textbox(s, cx + 0.20, yy + 0.13, cwid - 0.40, 0.26)
        para(tf, title, size=10.5, color=c, bold=True, first=True)
        tfb = textbox(s, cx + 0.20, yy + 0.44, cwid - 0.40, hgt - 0.52)
        for i, it in enumerate(items):
            rich(tfb, [("\u2022  ", {"color": c, "bold": True, "size": 10.5}),
                       (it, {"color": NAVY, "size": 10.5})],
                 first=(i == 0), space_after=3, line=1.18)
        yy += hgt + 0.18

    footer(s, n)
    notes(s, """
Briefly on perception. Adaptation handles what the robot can feel in its tracking error, not
something it has to look at. So the other thread is a vision language model on the robot's
own camera, streamed off the robot to a workstation running the model locally.

The benchmark shows a real tradeoff. Qwen at seven billion gives the best descriptions and
runs at a third of a hertz. Moondream is five times faster and the least accurate. Nothing in
this class fits inside a fifty hertz control loop, so this is a deliberative layer, not a
reflex.

To be clear about status: the output is still free text to a terminal. It is not connected to
the controller, and this talk does not claim more than that.
""")
    return s


def s_summary(prs, n):
    s = add_slide(prs)
    header(s, "What this summer produced", "Summary")

    stages = [
        ("Retargeting", ACCENT,
         "Fixed two retargeter defects: a joint-limit off-by-one and a 4.6 rad/s "
         "transient at the clip start."),
        ("Box pickup", ACCENT,
         "8.7 s whole-body pickup, carry and set-down. 33 training runs from "
         "first tracking to the deployed policy."),
        ("Hardware", ACCENT,
         "Exported to a numpy-only policy, deployed over ROS 2 at 50 Hz. "
         "57 logged runs on the physical X2."),
        ("Crawling", ACCENT,
         "Prone slope crawl with a purpose-written reward, contact and "
         "observation set. Also run on hardware."),
        ("Fault testing", RED,
         "Injected a right-knee actuator fault at 30% stiffness. Frozen policy "
         "falls from 14.7 s to 7.4 s."),
        ("Adaptation", GREEN,
         "Ported the ACC layer-adaptation law into Isaac and onto the robot. "
         "11.36 s, 6/6 set-down, p = 0.0011."),
        ("VLM", AMBER,
         "Head camera streamed to three local VLMs and benchmarked for latency "
         "and description quality."),
    ]
    nst = len(stages)
    gap = 0.20
    w = (CW - gap * (nst - 1)) / nst
    y, h = 1.80, 2.72
    bar = 0.42
    for i, (t, c, body) in enumerate(stages):
        x = ML + i * (w + gap)
        rrect(s, x, y, w, h, radius=0.055, fill=WHITE, line=RULE, lw=1.15)
        rect(s, x, y, w, bar, fill=c)
        tf = textbox(s, x + 0.06, y, w - 0.12, bar, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.CENTER)
        para(tf, t, size=11.5, color=WHITE, bold=True, first=True,
             align=PP_ALIGN.CENTER)
        tf = textbox(s, x + 0.14, y + bar + 0.16, w - 0.28, h - bar - 0.3)
        para(tf, body, size=10, color=NAVY, first=True, line=1.34)
        if i < nst - 1:
            arrow(s, x + w + 0.02, y + h / 2, x + w + gap - 0.02, y + h / 2,
                  color=RULE, lw=1.6)

    kicker(s, 4.94,
           "  the X2 integration and retargeting fixes, task and reward design in "
           "holosoma, every training run, the numpy export and ROS 2 deployment, "
           "the fault injection harness, and the adaptation port in both simulation "
           "and hardware form.",
           bold_prefix="Contributed here:", fill=LIGHT, size=12.5, h=0.84)

    tf = textbox(s, ML, 6.10, CW, 0.4)
    para(tf, "Built on holosoma for whole-body tracking, the OmniRetarget and CMU "
             "motion datasets, and the ACC 2026 adaptation law.",
         size=10.5, color=GRAY, first=True, italic=True, line=1.2)

    footer(s, n)
    notes(s, """
Stepping back, this is the arc. Retargeting first, which meant fixing the retargeter before
anything downstream worked. That fed the box pickup, thirty three training runs to reach the
policy that runs today, with fifty seven logged hardware runs behind it. Crawling was the
second behavior on the same infrastructure. Then the research question: a fault injection
harness, the frozen policy characterized, and the adaptation law ported into our environment.

On attribution. The tracking framework is holosoma, the motion data is public, and the
adaptation law is my mentor's. The contribution here is the X2 integration, the task and
reward design, the training, the hardware deployment, the fault experiments, and getting that
law working in an environment it had never run in.
""")
    return s


def s_next(prs, n):
    s = add_slide(prs)
    header(s, "Where the system stands, and what comes next",
           "Current status and next steps")

    # left: today
    lw = 5.90
    rrect(s, ML, BODY_TOP, lw, 5.10, radius=0.05, fill=RGBColor(0xF7, 0xFA, 0xFC),
          line=RULE, lw=1.1)
    tf = textbox(s, ML + 0.30, BODY_TOP + 0.24, lw - 0.6, 0.3)
    para(tf, "WHAT THE SYSTEM DOES TODAY", size=11, color=ACCENT, bold=True,
         first=True)
    items = [
        ("Box pickup", "8.7 s whole-body motion in simulation and on the "
                       "physical X2, with flat planted feet and a clean set-down."),
        ("Crawling", "Prone slope crawl trained in simulation, v3 exported and "
                     "run on hardware."),
        ("Deployment", "Numpy-only inference at 50 Hz over ROS 2, with a safety "
                       "ladder and per-tick logging."),
        ("Adaptation", "Runs online at 50 Hz in simulation and in the hardware "
                       "deploy loop, and recovers the task under a knee fault."),
    ]
    yy = BODY_TOP + 0.64
    for k, v in items:
        rect(s, ML + 0.30, yy + 0.04, 0.045, 0.86, fill=ACCENT)
        tf = textbox(s, ML + 0.50, yy, lw - 0.86, 0.92)
        para(tf, k, size=13.5, color=NAVY, bold=True, first=True, space_after=3)
        para(tf, v, size=11.5, color=GRAY, line=1.24)
        yy += 1.05

    # right: next
    x = ML + lw + 0.34
    w = CW - lw - 0.34
    rrect(s, x, BODY_TOP, w, 5.10, radius=0.05, fill=WHITE, line=RULE, lw=1.1)
    tf = textbox(s, x + 0.30, BODY_TOP + 0.24, w - 0.6, 0.3)
    para(tf, "NEXT STEPS, IN PRIORITY ORDER", size=11, color=GREEN, bold=True,
         first=True)
    nexts = [
        ("Generalize the fault result", "A hip or ankle fault, and a held-out "
         "seed pool."),
        ("Run the comparison on hardware", "Both arms already run on the robot. "
         "What is missing is degrading a physical joint."),
        ("A balance-targeted error signal", "Regulating joint error assumes "
         "balance follows. Define the error on a centre-of-mass term instead."),
        ("Gate the adaptation", "Engage above an error threshold to restore a "
         "no-worse-than-frozen floor."),
        ("Close the perception loop", "VLM output as a structured signal the "
         "controller can act on."),
    ]
    yy = BODY_TOP + 0.66
    for i, (k, v) in enumerate(nexts):
        tf = textbox(s, x + 0.30, yy, w - 0.60, 0.86)
        rich(tf, [(f"{i + 1}.  ", {"bold": True, "color": GREEN, "size": 13}),
                  (k, {"bold": True, "color": NAVY, "size": 13})],
             first=True, space_after=3, line=1.15)
        tf = textbox(s, x + 0.62, yy + 0.27, w - 0.92, 0.62)
        para(tf, v, size=11, color=GRAY, first=True, line=1.22)
        yy += 0.88

    footer(s, n)
    notes(s, """
Where things stand is on the left. Both behaviors exist, both have been on the physical
robot, and the adaptation runs online in simulation and in the hardware deploy loop.

On the right, what comes next. First, the honest limitation: this is one fault on one pool of
six seeds. A different joint and fresh seeds would make the claim properly.

Second, and this is the closest, running the comparison on the physical robot. Both arms
already run there. What is missing is degrading a real joint and alternating them enough
times to say something with confidence.

Third is a research direction. The adapter regulates joint error and assumes balance follows,
and the mask result is evidence that on a floating base it does not.
""")
    return s


def s_closing(prs, n):
    s = add_slide(prs)
    rect(s, 0, 0, SW, SH, fill=NAVY)
    rect(s, 0, 0, SW, 0.055, fill=ACCENT)

    tf = textbox(s, 1.20, 1.62, SW - 2.4, 2.6)
    para(tf, "Humanoid policies can learn manipulation and locomotion that works.",
         size=27, color=WHITE, bold=True, first=True, line=1.22)
    para(tf, "Resilience is a different problem. It requires the robot to respond "
             "to changes in its own dynamics.",
         size=27, color=RGBColor(0x7F, 0xB3, 0xD9), bold=True, space_before=14,
         line=1.22)

    rect(s, 1.20, 4.44, 1.42, 0.05, fill=ACCENT)

    tf = textbox(s, 1.20, 4.74, SW - 2.4, 0.9)
    para(tf, "This summer produced the behaviors, the deployment path, and the "
             "adaptation framework that gets us there, with evidence that "
             "adaptation recovers a task the frozen policy loses.",
         size=16, color=RGBColor(0xC8, 0xD6, 0xE4), first=True, line=1.34)

    tf = textbox(s, 1.20, 6.12, 6.5, 0.5)
    para(tf, "Baaqer Farhat  ·  Caltech ARCL  ·  SURF 2026", size=13,
         color=RGBColor(0x8E, 0xA3, 0xBA), first=True)
    tf = textbox(s, SW - 1.20 - 5.0, 6.12, 5.0, 0.5, align=PP_ALIGN.RIGHT)
    para(tf, "Thank you. Questions welcome.", size=13,
         color=RGBColor(0x8E, 0xA3, 0xBA), first=True, align=PP_ALIGN.RIGHT)

    notes(s, """
Getting a humanoid to pick up a box or crawl up a slope is a solvable problem, and it was
solved twice this summer, in simulation and on the real robot.

But a policy that works is not the same as a policy that is resilient. The moment the robot
stops being the robot the policy was trained on, and that happens constantly on hardware,
the behavior degrades and the policy cannot tell. Resilience means responding to changes in
your own dynamics.

Under a real actuator fault, adapting a single layer online, from nothing but tracking
error, recovers a task the frozen policy loses, on every seed tested.

Thank you. Happy to take questions.
""")
    return s


# ---------------------------------------------------------------- backup
def s_backup_divider(prs):
    s = add_slide(prs)
    rect(s, 0, 3.10, SW, 1.30, fill=LIGHT)
    tf = textbox(s, ML, 3.42, CW, 0.7, align=PP_ALIGN.CENTER)
    para(tf, "Backup", size=30, color=NAVY, bold=True, first=True,
         align=PP_ALIGN.CENTER)
    return s


def s_backup_healthy(prs):
    s = add_slide(prs)
    header(s, "The same law on a healthy robot spends margin the task needed",
           "Backup  ·  adaptation on a healthy policy")

    picture(s, ASSETS / "fig_healthy_vs_fault.png", ML, 1.52, w=7.05, h=3.10)

    x = ML + 7.35
    w = CW - 7.35
    bullets(s, x, 1.58, w, [
        ("The tracking claim reproduces. ",
         "On the window every variant survived, leg tracking error goes from "
         "14.82\u00b0 frozen to 8.01\u00b0 adapted. That is a 40% reduction, "
         "against 37% reported in the original MuJoCo testbed."),
        ("What does not transfer is the implication. ",
         "The published baseline falls at 1.94 s and never completes the motion, "
         "so adaptation could only add. Ours survives 14.68 s and completes the "
         "task on 6 of 6, and the same law spends that margin buying tracking "
         "accuracy the task did not need."),
        ("This is why the fault experiment matters. ",
         "A recovery controller applied to a healthy robot is outside the domain "
         "of the method. The comparison that means something is the one under the "
         "fault."),
    ], size=11.5, gap=10)

    kicker(s, 5.16,
           "  the numpy export matches the torch checkpoint to 4.2e-6 max action "
           "difference over 200 steps, and the baseline is that same numpy policy at "
           "gain zero, so the only difference between arms is the adaptation. "
           "A 1e-6 action perturbation moves leg error by 1.7\u00b0 over 2.4 s, so "
           "every comparison is multi-seed.",
           bold_prefix="Controls:", fill=WHITE, edge=RULE, size=11.5, h=0.86)

    tf = textbox(s, ML, 6.20, CW, 0.6)
    para(tf, "Inverse-inertia weighting (gx level 2) is worse in both regimes. "
             "That reproduces the original finding: the weighting aims the "
             "correction at the lightest joints, the wrists and head, whose "
             "tracking errors are irrelevant to the task.",
         size=11.5, color=GRAY, first=True, line=1.26)

    footer(s, 15)
    notes(s, """
This is the backup for the healthy robot claim.

The important honest point is the first bullet. The tracking result from the paper
reproduces cleanly in our environment. Leg tracking error drops forty percent, against
thirty seven percent in the original testbed. The law does what it says it does and it
transfers across simulators.

What does not transfer is the conclusion you would draw. In the original testbed the frozen
baseline falls after two seconds and never finishes the motion, so there was no balance
margin to lose and adaptation could only help. Our frozen baseline survives the full clip
and completes the task on every seed. So the same law is spending balance margin to buy
tracking accuracy that this task did not need, and the result is worse.

That is exactly why the fault experiment is the one that counts.
""")
    return s


def s_backup_s2r(prs):
    s = add_slide(prs)
    header(s, "Diagnosing a sim-to-real failure with matched simulation",
           "Backup  ·  v33 hardware")

    steps = [
        ("Symptom", RED,
         "All four takes of v33 abort on roll between 27% and 48% of the motion, "
         "always in the same direction."),
        ("First hypothesis, wrong", GRAY,
         "Compared against the only rollout on disk, which was a different policy. "
         "That pointed at a saturated ankle-roll command."),
        ("Matched baseline", ACCENT,
         "Identified the exact v33 reference clip and ran the same policy, same "
         "clip, 3 seeds. Simulation clips that joint on 92% of ticks and still "
         "succeeds, so saturation is learned behavior, not the gap."),
        ("Actual cause", GREEN,
         "The deployment fed the torso IMU into base_ang_vel, which holosoma "
         "builds from the pelvis. The torso sits three waist joints above it, so "
         "this substitutes both a different body and a different frame."),
        ("Confirmation", GREEN,
         "Reproducing that substitution in simulation drops the policy from 100% "
         "success and 9.8 s to 0% and 1.6 s, falling in the same 1.1 to 1.9 s "
         "window as the hardware."),
        ("Fix", GREEN,
         "A pelvis estimator composes the torso IMU with the measured waist "
         "joints. Removes 90.8% of the observation error and restores 100% "
         "success."),
    ]
    y = BODY_TOP + 0.04
    for t, c, body in steps:
        rect(s, ML, y + 0.05, 0.05, 0.68, fill=c)
        tf = textbox(s, ML + 0.24, y, 2.60, 0.7)
        para(tf, t, size=13, color=c, bold=True, first=True, line=1.16)
        tf = textbox(s, ML + 3.00, y, CW - 3.00, 0.72)
        para(tf, body, size=12, color=NAVY, first=True, line=1.24)
        y += 0.74

    kicker(s, y + 0.10,
           "  a single tracking-error number conflates command saturation, servo "
           "error and task error. Separating the three is what localized this, and "
           "reproducing the defect in simulation is what proved it.",
           bold_prefix="Method that found it:", fill=LIGHT, size=12.5, h=0.62)

    footer(s, 16)
    notes(s, """
This one is here in case it comes up, because the method matters more than the bug.

Version thirty three worked in simulation and failed on hardware every time, always
aborting on roll at roughly the same point in the motion.

The first diagnosis was wrong, and it was wrong for an instructive reason. The hardware logs
were compared against the only simulation rollout saved at the time, which came from a
different policy version. That mismatched baseline pointed at a saturated ankle joint.

Against a properly matched simulation, same policy, same clip, the saturation explanation
fell apart, because simulation saturates that joint even more than hardware does and still
succeeds.

The real cause was the deployment feeding the torso IMU into an observation the training
environment computes at the pelvis, and the torso sits three waist joints above the pelvis.
Wrong body, wrong frame.

The confirmation is what makes it conclusive. Reproducing that exact substitution inside
simulation, where the pelvis signal is available as ground truth, took the policy from one
hundred percent success to zero, falling in the same time window as the robot did. The
defect explains the failure on its own.
""")
    return s


def s_backup_config(prs):
    s = add_slide(prs)
    header(s, "Configuration and reproduction", "Backup  ·  details")

    cols = [
        ("TRAINING", ACCENT, [
            ("Framework", "holosoma whole-body tracking on IsaacLab / Isaac Sim"),
            ("Algorithm", "PPO, \u03bb = 0.95, \u03b3 = 0.99, 4096 envs"),
            ("Rates", "200 Hz physics, 50 Hz control"),
            ("Actor", "blind MLP 512-256-128 (ELU), 31 joint targets"),
            ("Critic", "privileged, adds base linear velocity and object state"),
            ("Pickup DR", "box mass 0.4 to 1.6 kg, friction, PD gain, action "
                          "delay, encoder noise, pushes"),
            ("Motion", "OmniRetarget sub3_largebox_003, CMU crawl_111_03"),
        ]),
        ("ADAPTATION", GREEN, [
            ("Law", "ACC 2026 causality-guided layer adaptation, W\u2080 leak"),
            ("Adapted", "layer index 2, 256 \u00d7 128 = 32,768 weights"),
            ("Gain", "\u0393 = 3e-4, leak 0.01, engage step 0"),
            ("Mask", "waist yaw, pitch, roll"),
            ("Fault", "right_knee_joint PD stiffness \u00d7 0.3"),
            ("Protocol", "6 seeds, 734 steps, no DR, no observation noise"),
            ("Baseline", "same numpy policy at gain 0"),
        ]),
    ]
    w = (CW - 0.42) / 2
    for ci, (title, c, rows) in enumerate(cols):
        x = ML + ci * (w + 0.42)
        rrect(s, x, BODY_TOP, w, 4.30, radius=0.05, fill=WHITE, line=RULE, lw=1.1)
        rect(s, x, BODY_TOP, w, 0.40, fill=c)
        tf = textbox(s, x + 0.22, BODY_TOP, w - 0.44, 0.40, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, title, size=11, color=WHITE, bold=True, first=True)
        yy = BODY_TOP + 0.58
        for k, v in rows:
            tf = textbox(s, x + 0.24, yy, 1.62, 0.5)
            para(tf, k, size=11, color=ACCENT if ci == 0 else GREEN, bold=True,
                 first=True)
            tf = textbox(s, x + 1.94, yy, w - 2.18, 0.52)
            para(tf, v, size=11, color=NAVY, first=True, line=1.2)
            yy += 0.525

    rrect(s, ML, 5.94, CW, 0.78, radius=0.05, fill=RGBColor(0xF7, 0xFA, 0xFC),
          line=RULE, lw=1.1)
    tf = textbox(s, ML + 0.26, 6.06, CW - 0.52, 0.56)
    para(tf, "Full reproduction instructions, per-seed results and the hardware "
             "protocol are in adaptation/README.md and SETUP_ISAAC.md.",
         size=11.5, color=NAVY, first=True, space_after=4)
    para(tf, "Repository: Agibot-humanoid  ·  box_pickup/  adaptation/  "
             "agibot_control_functions/", size=10, color=GRAY, font=MONO)

    footer(s, 17)
    notes(s, "Configuration reference. Use if asked about specific "
             "hyperparameters, the fault definition, or how to reproduce a run.")
    return s


# ================================================================== main
def main():
    prs = new_deck()
    builders = [s_title, s_problem, s_system, s_box_task, s_box_lessons,
                s_box_diagnosis, s_crawl, s_why_adapt, s_method, s_results,
                s_adapt_videos, s_vlm, s_summary, s_next, s_closing]
    for i, b in enumerate(builders, start=1):
        b(prs, i)
    s_backup_divider(prs)
    s_backup_healthy(prs)
    s_backup_s2r(prs)
    s_backup_config(prs)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
