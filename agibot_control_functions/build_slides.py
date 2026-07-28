#!/usr/bin/env python3
"""Build the progress-update deck (PPTX + PDF).

Topics: (1) box-pickup motion-retargeting policy, (2) 3-VLM perception setup.
LaTeX math is rendered to PNG via matplotlib mathtext, videos are embedded, and
the deck is exported to PDF via PowerPoint COM automation.
"""
import os
import subprocess

import cv2
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))          # Agibot-humanoid
VID = os.path.join(ROOT, "box_pickup", "videos")
FRAME224 = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "captures", "my_run", "frame_00224.jpg")
OUTDIR = os.path.join(ROOT, "slides")
ASSETS = os.path.join(OUTDIR, "assets")
os.makedirs(ASSETS, exist_ok=True)

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
GRAY = RGBColor(0x44, 0x4A, 0x54)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x12, 0x8A, 0x5E)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def eq_png(latex, name, fontsize=22):
    path = os.path.join(ASSETS, name)
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.text(0, 0, f"${latex}$", fontsize=fontsize, color="#0B2A4A")
    fig.savefig(path, dpi=200, bbox_inches="tight", pad_inches=0.08,
                transparent=True)
    plt.close(fig)
    return path


def poster(video_path, name):
    path = os.path.join(ASSETS, name)
    cap = cv2.VideoCapture(video_path)
    n = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
    cap.set(cv2.CAP_PROP_POS_FRAMES, max(0, n // 2))
    ok, frame = cap.read()
    if not ok:
        cap.set(cv2.CAP_PROP_POS_FRAMES, 0)
        ok, frame = cap.read()
    cap.release()
    cv2.imwrite(path, frame)
    return path, frame.shape[1], frame.shape[0]


def img_size(path):
    im = cv2.imread(path)
    return im.shape[1], im.shape[0]


# ---------- slide helpers ----------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rrect(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def center_text(shape, lines, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), text, size, color, bold=bold)


def pipeline(slide, stages, x, y, w, h, box_color=BLUE, txt=WHITE):
    n = len(stages)
    gap = Inches(0.42)
    bw = int((w - gap * (n - 1)) / n)
    cx = x
    for i, st in enumerate(stages):
        box = rrect(slide, cx, y, bw, h, box_color)
        lines = [(ln, 12, txt, True) for ln in st.split("\n")]
        center_text(box, lines)
        if i < n - 1:
            ah = Inches(0.34)
            ar = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, cx + bw + Emu(int(gap * 0.12)),
                y + int((h - ah) / 2), Emu(int(gap * 0.76)), ah)
            ar.fill.solid()
            ar.fill.fore_color.rgb = ACCENT
            ar.line.fill.background()
            ar.shadow.inherit = False
        cx += bw + gap


def stat_card(slide, x, y, w, h, number, label):
    box = rrect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, w, Inches(0.08), ACCENT)
    center_text(box, [(number, 26, NAVY, True), (label, 11, GRAY, False)])
    return box


def placeholder(slide, x, y, w, h, caption):
    box = rrect(slide, x, y, w, h, RGBColor(0xED, 0xF1, 0xF6),
                line=RGBColor(0xB6, 0xC6, 0xD6))
    center_text(box, [("Figure placeholder", 15, BLUE, True),
                      (caption, 12, GRAY, False)])
    return box


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color=GRAY, bold=False, italic=False,
            font="Segoe UI"):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def header(slide, title, kicker=None):
    rect(slide, 0, 0, EMU_W, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), EMU_W, Inches(0.06), ACCENT)
    tf = textbox(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.95),
                 MSO_ANCHOR.MIDDLE)
    if kicker:
        p = tf.paragraphs[0]
        set_run(p.add_run(), kicker, 12, RGBColor(0x9F, 0xC7, 0xE8), bold=True)
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), title, 26, WHITE, bold=True)
    else:
        set_run(tf.paragraphs[0].add_run(), title, 26, WHITE, bold=True)


def bullets(slide, x, y, w, h, items, size=16, gap=6):
    tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        level = 0
        text = it
        if isinstance(it, tuple):
            text, level = it
        p.level = level
        run = p.add_run()
        bullet = "•  " if level == 0 else "–  "
        set_run(run, bullet + text, size - level * 1,
                GRAY if level else NAVY, bold=(level == 0 and False))
    return tf


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H

    # ---- assets ----
    eq_retarget = eq_png(
        r"\min_{\Delta q}\;\|L(q)-L_{human}\|^2+w_s\|\Delta q\|^2"
        r"\quad s.t.\;\;\|\Delta q\|\leq\delta,\;\; q^-\!\leq q\leq q^+",
        "eq_retarget.png", 20)
    eq_reward = eq_png(
        r"r_{hands}=\exp\!\left(-\frac{\sum_{i\in\{L,R\}}"
        r"\|p^{hand}_i-p^{box}\|^2}{\sigma^2}\right),\;\;\sigma=0.25\,m",
        "eq_reward.png", 20)
    p_human, hw, hh = poster(
        os.path.join(VID, "omomo_human_mocap_sub3_largebox_003.mp4"),
        "poster_human.png")
    p_x2, xw, xh = poster(
        os.path.join(VID, "reference_motion_FIXED_sub3_largebox_003.mp4"),
        "poster_x2.png")

    # ================= Slide 1 : Title =================
    s = blank(prs)
    rect(s, 0, 0, EMU_W, EMU_H, NAVY)
    rect(s, 0, Inches(4.05), EMU_W, Inches(0.08), ACCENT)
    tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.8))
    set_run(tf.paragraphs[0].add_run(), "AgiBot X2 — Progress Update", 40,
            WHITE, bold=True)
    p = tf.add_paragraph()
    set_run(p.add_run(),
            "Box-Pickup Motion-Retargeting Policy  |  On-Robot VLM Perception",
            20, RGBColor(0x9F, 0xC7, 0xE8))
    tf2 = textbox(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(1.0))
    set_run(tf2.paragraphs[0].add_run(),
            "Whole-body manipulation via human MoCap retargeting + "
            "vision-language perception benchmark", 15,
            RGBColor(0xC9, 0xD6, 0xE3))

    # ========= Slide 2 : Policy objective & approach =========
    s = blank(prs)
    header(s, "Box-Pickup Policy — Objective & Approach", "NEW POLICY")
    bullets(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(3.2), [
        "Task: whole-body box pickup — walk up, two-hand grasp, lift, "
        "carry, and set down a large box.",
        "Reference source: human motion-capture clip "
        "(OmniRetarget dataset, sub3_largebox_003, SMPL-H 52-joint model).",
        "Embodiment switch: motion authored on a human SMPL-H skeleton is "
        "retargeted onto the AgiBot X2 — 31 DOF (6+6 legs, 3 waist, 7+7 "
        "arms, 2 head), 1.4 m tall, rigid half-sphere palms (no fingers).",
        "Status: policy is currently in training; deployment on the physical "
        "X2 is planned once this run completes.",
    ], size=16, gap=9)
    # visual pipeline
    tf = textbox(s, Inches(0.5), Inches(4.55), Inches(6), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(), "Pipeline", 15, BLUE, bold=True)
    pipeline(s, [
        "Human SMPL-H\nMoCap",
        "Interaction-Mesh\nRetargeting",
        "X2 31-DOF\nReference Traj.",
        "RL Tracking Policy\n(Isaac Sim / PPO)",
        "Deploy on\nphysical X2",
    ], Inches(0.5), Inches(5.05), Inches(12.3), Inches(1.15))
    # mark final stage as "next" (accent outline via lighter color)
    tf = textbox(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(),
            "Blind whole-body policy (no cameras); object state used only by "
            "the critic during training.", 12, GRAY, italic=True)

    # ========= Slide 3 : Retargeting logic =========
    s = blank(prs)
    header(s, "Retargeting Logic", "NEW POLICY")
    # left: text bullets (top)
    bullets(s, Inches(0.5), Inches(1.5), Inches(7.0), Inches(3.2), [
        "Method: interaction-mesh retargeting solved as a per-frame convex "
        "QP with SQP linearization (CVXPY + Clarabel).",
        "14-landmark SMPL-H → X2 correspondence (pelvis, hips, knees, "
        "ankles, toes, shoulders, elbows, wrists → hand-contact links).",
        "Synthetic hand-contact link at (0.02, 0, -0.13) m from each wrist "
        "so the human wrist maps to the grasp surface.",
        "Constraints: Laplacian (interaction-mesh) preservation; "
        "non-penetration (floor + box); foot-stick / foot-lock; joint "
        "limits; trust region; self-collision.",
    ], size=14, gap=7)
    # left-bottom: formula box under the text
    rrect(s, Inches(0.5), Inches(5.05), Inches(7.0), Inches(1.9), LIGHT)
    tf = textbox(s, Inches(0.75), Inches(5.15), Inches(6.5), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(), "Per-frame QP objective", 13,
            BLUE, bold=True)
    ew, eh = img_size(eq_retarget)
    w = Inches(6.3)
    s.shapes.add_picture(eq_retarget, Inches(0.85), Inches(5.7), width=w)
    tf = textbox(s, Inches(0.75), Inches(6.5), Inches(6.5), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(),
            "L(·): interaction-mesh Laplacian · δ: trust-region step · "
            "q±: joint limits", 10.5, GRAY, italic=True)
    # right: space for a figure from the paper
    placeholder(s, Inches(7.75), Inches(1.5), Inches(5.05), Inches(5.45),
                "Add interaction-mesh / retargeting figure from the paper")

    # ========= Slide 4 : reference videos =========
    s = blank(prs)
    header(s, "Reference Motion:  Human MoCap  →  Retargeted X2", "NEW POLICY")
    vy = Inches(1.6)
    vw = Inches(5.9)
    vh1 = Emu(int(vw * hh / hw))
    vh2 = Emu(int(vw * xh / xw))
    # left = human
    tf = textbox(s, Inches(0.5), Inches(1.25), vw, Inches(0.35))
    set_run(tf.paragraphs[0].add_run(),
            "Source — human SMPL-H MoCap", 15, NAVY, bold=True)
    s.shapes.add_movie(
        os.path.join(VID, "omomo_human_mocap_sub3_largebox_003.mp4"),
        Inches(0.5), vy, vw, vh1, poster_frame_image=p_human,
        mime_type="video/mp4")
    # right = retargeted x2
    tf = textbox(s, Inches(6.9), Inches(1.25), vw, Inches(0.35))
    set_run(tf.paragraphs[0].add_run(),
            "Retargeted — AgiBot X2 (31-DOF)", 15, NAVY, bold=True)
    s.shapes.add_movie(
        os.path.join(VID, "reference_motion_FIXED_sub3_largebox_003.mp4"),
        Inches(6.9), vy, vw, vh2, poster_frame_image=p_x2,
        mime_type="video/mp4")
    tf = textbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6))
    set_run(tf.paragraphs[0].add_run(),
            "Same clip before and after the interaction-mesh IK pipeline; "
            "contact states corrected post-retargeting. (Click to play.)",
            12, GRAY, italic=True)

    # ========= Slide 5 : training =========
    s = blank(prs)
    header(s, "Policy Training — RL Tracking", "NEW POLICY")
    # stat cards row
    cards = [("4096", "parallel envs"), ("50 Hz", "control (200 Hz sim)"),
             ("31", "actuated DOF"), ("PPO", "λ=0.95 · γ=0.99")]
    cx = Inches(0.5)
    cw = Inches(2.9)
    cgap = Inches(0.23)
    for num, lab in cards:
        stat_card(s, cx, Inches(1.5), cw, Inches(1.15), num, lab)
        cx += cw + cgap
    # left: bullets
    bullets(s, Inches(0.5), Inches(2.95), Inches(7.0), Inches(4.0), [
        "Framework: IsaacLab (Isaac Sim) + holosoma.",
        "Actor: blind MLP 512-256-128 (ELU) — prev action, base ang-vel, "
        "dof pos/vel, motion phase, ref orientation.",
        "Critic: privileged obs adds base lin-vel + object & body state.",
        "Key reward: two-hand-to-box proximity as a product term — both "
        "palms must reach the box simultaneously.",
        "Status: currently training; deploy on physical X2 after the run.",
    ], size=14, gap=8)
    # right: actor / critic mini-diagram
    rrect(s, Inches(7.75), Inches(2.95), Inches(5.05), Inches(0.75), BLUE)
    center_text(s.shapes[-1], [("Actor  (blind)  →  512-256-128  →  31 DOF",
                                13, WHITE, True)])
    rrect(s, Inches(7.75), Inches(3.85), Inches(5.05), Inches(0.75), NAVY)
    center_text(s.shapes[-1], [("Critic  (privileged: + object & body state)",
                                13, WHITE, True)])
    # right: reward formula box
    rrect(s, Inches(7.75), Inches(4.85), Inches(5.05), Inches(2.1), LIGHT)
    tf = textbox(s, Inches(7.95), Inches(4.95), Inches(4.6), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(), "Two-hand grasp reward", 13,
            BLUE, bold=True)
    s.shapes.add_picture(eq_reward, Inches(7.95), Inches(5.5),
                         width=Inches(4.6))
    tf = textbox(s, Inches(7.95), Inches(6.45), Inches(4.6), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(),
            "Product (not mean) over hands forces bimanual contact.",
            11, GRAY, italic=True)

    # ========= Slide 6 : VLM setup =========
    s = blank(prs)
    header(s, "VLM Perception — Setup", "PERCEPTION")
    bullets(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0), [
        "Goal: scene understanding from the X2 head RGB camera; frames "
        "streamed to a workstation for evaluation.",
        "Three VLMs served locally via Ollama: "
        "Qwen2.5-VL 7B, Qwen2.5-VL 3B, Moondream (~1.8B).",
        "Protocol: 15 s clip @ ~20 fps → 12 evenly-spaced frames; identical "
        "prompt, temperature 0.2; one model resident at a time (8 GB VRAM).",
    ], size=16, gap=8)
    rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.4), LIGHT)
    tf = textbox(s, Inches(0.8), Inches(4.65), Inches(11.7), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(), "Prompt (identical for all models)",
            14, BLUE, bold=True)
    tf = textbox(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.6))
    set_run(tf.paragraphs[0].add_run(),
            "\u201cYou are a robot's vision system. List the main objects and "
            "people you see and roughly where they are (left/center/right). "
            "Then one short summary sentence. Be concise.\u201d",
            16, NAVY, italic=True)

    # ========= Slide 7 : metrics =========
    s = blank(prs)
    header(s, "VLM Performance Metrics", "PERCEPTION")
    rows = [
        ("Model", "Params", "Mean lat.", "Median", "Rate", "Gen speed"),
        ("Qwen2.5-VL 7B", "7 B", "2.87 s", "2.89 s", "0.35 Hz", "51.3 tok/s"),
        ("Qwen2.5-VL 3B", "3 B", "2.45 s", "2.50 s", "0.41 Hz", "107.6 tok/s"),
        ("Moondream", "~1.8 B", "0.61 s", "0.62 s", "1.65 Hz", "215.4 tok/s"),
    ]
    tbl_x, tbl_y = Inches(0.7), Inches(1.7)
    tbl_w, tbl_h = Inches(11.9), Inches(2.6)
    gtbl = s.shapes.add_table(len(rows), 6, tbl_x, tbl_y, tbl_w, tbl_h).table
    widths = [3.2, 1.6, 1.9, 1.7, 1.7, 1.8]
    for c, wv in enumerate(widths):
        gtbl.columns[c].width = Inches(wv)
    for r_i, row in enumerate(rows):
        for c_i, val in enumerate(row):
            cell = gtbl.cell(r_i, c_i)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.CENTER
            run = para.add_run()
            if r_i == 0:
                set_run(run, val, 15, WHITE, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                highlight = (row[0] == "Moondream")
                set_run(run, val, 14,
                        ACCENT if highlight and c_i >= 4 else GRAY,
                        bold=(c_i == 0 or (highlight and c_i >= 4)))
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r_i % 2 else WHITE
    bullets(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(2.4), [
        "Warm, end-to-end latencies on an RTX 5070 (8 GB); 12 frames, "
        "single 15 s clip.",
        "Moondream is ~4-5× faster — the only model near real-time (~1.7 Hz).",
        "Qwen2.5-VL 7B gives the most accurate scene descriptions; 3B is "
        "slower than Moondream yet less accurate on this scene.",
    ], size=15, gap=7)

    # ========= Slide 8 : qualitative frame 224 =========
    s = blank(prs)
    header(s, "VLM Qualitative Comparison — Frame 224 (t = 11.3 s)",
           "PERCEPTION")
    iw, ih = img_size(FRAME224)
    pic_w = Inches(4.6)
    pic_h = Emu(int(pic_w * ih / iw))
    s.shapes.add_picture(FRAME224, Inches(0.5), Inches(1.6), width=pic_w)
    outs = [
        ("Qwen2.5-VL 7B  (2.80 s | 51.6 tok/s)",
         "Objects: yellow toolbox, green cord, wheels, white t-shirt with "
         "text, gray pants. People: man in white t-shirt squatting. "
         "Summary: a person interacting with a wheeled device in an "
         "industrial setting."),
        ("Qwen2.5-VL 3B  (2.60 s | 107.9 tok/s)",
         "A person squatting in the center of the room; yellow and black "
         "robotic feet on the floor, likely attached to a robot or machine."),
        ("Moondream  (0.67 s | 203.4 tok/s)",
         "Person in white shirt, gray pants, black shoes, and yellow shoes "
         "standing near a desk with a green cord on it."),
    ]
    y = Inches(1.6)
    bx = Inches(5.4)
    bw = Inches(7.4)
    bh = Inches(1.65)
    for title, body in outs:
        rect(s, bx, y, bw, bh, LIGHT)
        rect(s, bx, y, Inches(0.09), bh, ACCENT)
        tf = textbox(s, bx + Inches(0.25), y + Inches(0.08), bw - Inches(0.4),
                     bh - Inches(0.16))
        set_run(tf.paragraphs[0].add_run(), title, 14, BLUE, bold=True)
        p = tf.add_paragraph()
        set_run(p.add_run(), body, 12.5, GRAY)
        y = y + bh + Inches(0.18)

    pptx_path = os.path.join(OUTDIR, "AgiBot_X2_Progress_Update.pptx")
    prs.save(pptx_path)
    print("Saved", pptx_path)
    return pptx_path


def to_pdf(pptx_path):
    pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
    ps = (
        "$p=New-Object -ComObject PowerPoint.Application;"
        f"$d=$p.Presentations.Open('{pptx_path}',$true,$false,$false);"
        f"$d.SaveAs('{pdf_path}',32);$d.Close();$p.Quit()"
    )
    subprocess.run(["powershell", "-NoProfile", "-Command", ps],
                   capture_output=True, timeout=180)
    print("Saved", pdf_path if os.path.isfile(pdf_path) else "(PDF FAILED)")


if __name__ == "__main__":
    p = build()
    to_pdf(p)
