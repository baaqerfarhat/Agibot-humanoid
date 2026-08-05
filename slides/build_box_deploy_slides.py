#!/usr/bin/env python3
"""Build X2 box-pickup deployment analysis slides from run_logs canvas data.

Style matches presentation/build_slides.py (navy band, accent rule, hanging bullets).
Data source: run_logs/_run_analysis.json (same series as the deployment canvas).
Output: slides/X2_Box_Pickup_Deployment.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
DATA = ROOT / "run_logs" / "_run_analysis.json"
OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "deploy_assets"

NAVY = RGBColor(0x15, 0x22, 0x40)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
DARK = RGBColor(0x23, 0x28, 0x2E)
GRAY = RGBColor(0x5C, 0x66, 0x70)
LIGHT = RGBColor(0xEF, 0xF3, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = "#C0392B"
BLUE = "#2E86C1"
MUTED = "#8A97A6"
FONT = "Arial"
EQ_HEX = "#20242B"

IN = 914400
SLIDE_W = Emu(int(13.333 * IN))
SLIDE_H = Emu(int(7.5 * IN))


def emu(inches: float) -> Emu:
    return Emu(int(inches * IN))


# ----------------------------- plots -----------------------------
def _style_axes(ax, title: str, xlabel: str, ylabel: str):
    ax.set_title(title, fontsize=12, fontweight="bold", color="#152240", pad=8)
    ax.set_xlabel(xlabel, fontsize=10, color="#5C6670")
    ax.set_ylabel(ylabel, fontsize=10, color="#5C6670")
    ax.tick_params(colors="#5C6670", labelsize=9)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.grid(True, axis="y", color="#E6EBF0", linewidth=0.8)
    ax.set_facecolor("white")


def _shade_phases(ax, t, phases):
    colors = {
        "ramp": "#E8ECF0",
        "settle": "#D6EAF8",
        "pickup": "#D4E6F1",
        "policy": "#D4E6F1",
        "carry": "#D5F5E3",
        "setdown": "#FCF3CF",
        "done": "#E8DAEF",
    }
    if not phases:
        return
    start = 0
    cur = phases[0]
    for i in range(1, len(phases) + 1):
        if i == len(phases) or phases[i] != cur:
            ax.axvspan(t[start], t[min(i, len(t) - 1)], color=colors.get(cur, "#F0F0F0"),
                       alpha=0.55, lw=0)
            if i < len(phases):
                start = i
                cur = phases[i]


def _run_by_ts(data, ts: str):
    for r in data["runs"]:
        if r["timestamp"] == ts:
            return r
    raise KeyError(ts)


def make_plots(data: dict) -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}

    # --- 1. foot-edge by run ---
    runs = data["runs"]
    labels = []
    edges = []
    colors = []
    for r in runs:
        tag = "v27" if "v27" in r["policy"] else "v30"
        labels.append(f"{tag}\n{r['timestamp'][9:13]}")
        edges.append(r["metrics"]["foot_edge_score"])
        colors.append("#7F8C8D" if tag == "v27" else "#C0392B")
    fig, ax = plt.subplots(figsize=(11.5, 3.6), dpi=160)
    ax.bar(range(len(edges)), edges, color=colors, width=0.72)
    ax.axhline(3.0, color="#1F8A65", ls="--", lw=1.6, label="flat-foot target (3°)")
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=7.5)
    _style_axes(ax, "Foot-edge score by hardware run (RMS ankle-roll deviation, pickup window)",
                "Run", "Foot-edge (°)")
    ax.legend(frameon=False, fontsize=9, loc="upper left")
    fig.tight_layout()
    p = ASSET_DIR / "foot_edge_by_run.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["foot_edge"] = p

    # --- 2. hybrid desired vs actual ankles + base ---
    hybrid = _run_by_ts(data, "20260729_173741")
    tr = hybrid["traces"]
    t = np.asarray(tr["t"], dtype=float)
    phases = tr.get("phase", [])

    fig, axes = plt.subplots(2, 2, figsize=(11.8, 5.6), dpi=160, sharex=True)
    pairs = [
        (axes[0, 0], "right_ankle_roll_joint", "Right ankle-roll (foot-edge)"),
        (axes[0, 1], "left_ankle_roll_joint", "Left ankle-roll (foot-edge)"),
        (axes[1, 0], "right_knee_joint", "Right knee"),
        (axes[1, 1], "left_knee_joint", "Left knee"),
    ]
    for ax, j, title in pairs:
        _shade_phases(ax, t, phases)
        ax.plot(t, tr[f"{j}__tgt"], color=MUTED, lw=1.8, label="Commanded")
        ax.plot(t, tr[f"{j}__meas"], color=RED, lw=1.6, label="Measured")
        _style_axes(ax, title, "Time (s)", "Angle (°)")
        ax.legend(frameon=False, fontsize=8, loc="upper right")
    fig.suptitle("Hybrid run 17:37:41  ·  commanded target vs measured encoder",
                 fontsize=13, fontweight="bold", color="#152240", y=1.01)
    fig.tight_layout()
    p = ASSET_DIR / "hybrid_joints_ankles_knees.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["hybrid_joints"] = p

    # hips + ankles pitch
    fig, axes = plt.subplots(2, 2, figsize=(11.8, 5.6), dpi=160, sharex=True)
    pairs = [
        (axes[0, 0], "right_hip_pitch_joint", "Right hip-pitch"),
        (axes[0, 1], "left_hip_pitch_joint", "Left hip-pitch"),
        (axes[1, 0], "right_ankle_pitch_joint", "Right ankle-pitch"),
        (axes[1, 1], "left_ankle_pitch_joint", "Left ankle-pitch"),
    ]
    for ax, j, title in pairs:
        _shade_phases(ax, t, phases)
        ax.plot(t, tr[f"{j}__tgt"], color=MUTED, lw=1.8, label="Commanded")
        ax.plot(t, tr[f"{j}__meas"], color=RED, lw=1.6, label="Measured")
        _style_axes(ax, title, "Time (s)", "Angle (°)")
        ax.legend(frameon=False, fontsize=8, loc="best")
    fig.suptitle("Hybrid run 17:37:41  ·  hip / ankle-pitch tracking",
                 fontsize=13, fontweight="bold", color="#152240", y=1.01)
    fig.tight_layout()
    p = ASSET_DIR / "hybrid_joints_hips.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["hybrid_hips"] = p

    # base attitude
    fig, ax = plt.subplots(figsize=(11.5, 3.8), dpi=160)
    _shade_phases(ax, t, phases)
    ax.plot(t, tr["roll"], color=RED, lw=1.7, label="Torso roll")
    ax.plot(t, tr["pitch"], color=BLUE, lw=1.7, label="Torso pitch")
    _style_axes(ax, "Base attitude during hybrid run (fixation wobble)",
                "Time (s)", "Angle (°)")
    ax.legend(frameon=False, fontsize=9)
    fig.tight_layout()
    p = ASSET_DIR / "hybrid_base_attitude.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["hybrid_base"] = p

    # worst pickup ankle edge
    worst = _run_by_ts(data, "20260729_175222")
    tw = np.asarray(worst["traces"]["t"], dtype=float)
    pw = worst["traces"].get("phase", [])
    fig, axes = plt.subplots(1, 2, figsize=(11.5, 3.6), dpi=160, sharey=True)
    for ax, j, title in [
        (axes[0], "right_ankle_roll_joint", "Right ankle-roll — worst run (13.1°)"),
        (axes[1], "left_ankle_roll_joint", "Left ankle-roll — worst run"),
    ]:
        _shade_phases(ax, tw, pw)
        ax.plot(tw, worst["traces"][f"{j}__tgt"], color=MUTED, lw=1.8, label="Commanded")
        ax.plot(tw, worst["traces"][f"{j}__meas"], color=RED, lw=1.6, label="Measured")
        _style_axes(ax, title, "Time (s)", "Angle (°)")
        ax.legend(frameon=False, fontsize=8)
    fig.tight_layout()
    p = ASSET_DIR / "worst_ankle_roll.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["worst_ankle"] = p

    # per-joint RMSE for hybrid
    order = [
        ("left_hip_pitch_joint", "L hip-pit"),
        ("left_hip_roll_joint", "L hip-rol"),
        ("left_hip_yaw_joint", "L hip-yaw"),
        ("left_knee_joint", "L knee"),
        ("left_ankle_pitch_joint", "L ank-pit"),
        ("left_ankle_roll_joint", "L ank-rol"),
        ("right_hip_pitch_joint", "R hip-pit"),
        ("right_hip_roll_joint", "R hip-rol"),
        ("right_hip_yaw_joint", "R hip-yaw"),
        ("right_knee_joint", "R knee"),
        ("right_ankle_pitch_joint", "R ank-pit"),
        ("right_ankle_roll_joint", "R ank-rol"),
    ]
    vals = [hybrid["metrics"]["per_joint_rmse"][k] for k, _ in order]
    labs = [s for _, s in order]
    fig, ax = plt.subplots(figsize=(11.5, 3.8), dpi=160)
    bar_c = ["#C0392B" if "ank-rol" in s else BLUE for s in labs]
    ax.barh(range(len(vals)), vals, color=bar_c, height=0.7)
    ax.set_yticks(range(len(labs)))
    ax.set_yticklabels(labs, fontsize=9)
    ax.invert_yaxis()
    _style_axes(ax, "Leg tracking RMSE by joint — hybrid run 17:37:41 (active phases)",
                "RMSE (°)", "")
    fig.tight_layout()
    p = ASSET_DIR / "per_joint_rmse.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["per_joint"] = p

    # metric summary cards as a small table image
    v30 = [r for r in runs if "v30" in r["policy"]]
    v27 = [r for r in runs if "v27" in r["policy"]]

    def med(xs):
        xs = sorted(xs)
        return xs[len(xs) // 2]

    fig, ax = plt.subplots(figsize=(10, 2.8), dpi=160)
    ax.axis("off")
    cell = [
        ["", "v27 (n=4)", "v30 deployed (n=15)", "Target"],
        ["Median foot-edge (°)", f"{med([r['metrics']['foot_edge_score'] for r in v27]):.1f}",
         f"{med([r['metrics']['foot_edge_score'] for r in v30]):.1f}", "< 3"],
        ["Median roll σ (°)", f"{med([r['metrics']['base_roll_std'] for r in v27]):.1f}",
         f"{med([r['metrics']['base_roll_std'] for r in v30]):.1f}", "low"],
        ["Worst foot-edge (°)", f"{max(r['metrics']['foot_edge_score'] for r in v27):.1f}",
         f"{max(r['metrics']['foot_edge_score'] for r in v30):.1f}", "—"],
        ["Median leg RMSE (°)", f"{med([r['metrics']['leg_track_rmse_active'] for r in v27]):.1f}",
         f"{med([r['metrics']['leg_track_rmse_active'] for r in v30]):.1f}", "—"],
    ]
    table = ax.table(cellText=cell, loc="center", cellLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(11)
    table.scale(1.2, 1.55)
    for (r, c), cell_obj in table.get_celld().items():
        cell_obj.set_edgecolor("#D5DDE5")
        if r == 0:
            cell_obj.set_facecolor("#152240")
            cell_obj.get_text().set_color("white")
            cell_obj.get_text().set_fontweight("bold")
        elif c == 0:
            cell_obj.set_facecolor("#EFF3F7")
            cell_obj.get_text().set_fontweight("bold")
            cell_obj.get_text().set_color("#2E86C1")
        else:
            cell_obj.set_facecolor("white")
            cell_obj.get_text().set_color("#23282E")
    fig.tight_layout()
    p = ASSET_DIR / "summary_table.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    paths["summary"] = p

    # penalty equations
    eqs = {
        "foot_slip": r"$r_{\mathrm{slip}}=\sum_f \|v^{xy}_f\|\,\mathbf{1}[\|F_f\|>\tau]$",
        "contact_loss": r"$r_{\mathrm{air}}=\sum_f \mathbf{1}[\|F_f\|<\tau]$",
        "anchor": r"$r_{\mathrm{anchor}}=\sum_f \|p^{xy}_f-p^{xy}_{f,0}\|$",
        "not_flat": r"$r_{\mathrm{tilt}}=\sum_f (w_r|g^y_f|+w_p|g^x_f|)\,\mathbf{1}_{\mathrm{contact}}$",
        "edge": r"$r_{\mathrm{edge}}=\sum_f \mathbf{1}[\mathrm{contact}\wedge\sin\theta_f>\theta_{\mathrm{th}}]$",
    }
    for name, tex in eqs.items():
        fig = plt.figure(figsize=(7.5, 0.85))
        fig.text(0.0, 0.5, tex, va="center", ha="left", fontsize=22, color=EQ_HEX)
        out = ASSET_DIR / f"eq_{name}.png"
        fig.savefig(out, dpi=200, transparent=True, bbox_inches="tight", pad_inches=0.04)
        plt.close(fig)
        paths[f"eq_{name}"] = out

    return paths


# ----------------------------- pptx primitives -----------------------------
def add_bg(slide, color):
    s = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)


def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=8, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, bold, color) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
    return tb


def title_band(slide, title):
    add_rect(slide, 0, 0, SLIDE_W, emu(1.1), NAVY)
    add_rect(slide, 0, emu(1.1), SLIDE_W, emu(0.055), ACCENT)
    add_text(slide, emu(0.7), emu(0.24), emu(12), emu(0.62),
             [[(title, 26, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n):
    add_text(slide, emu(0.7), emu(7.04), emu(10), emu(0.36),
             [[("X2 box-pickup deployment analysis  ·  hardware run_logs", 9, False, GRAY)]])
    add_text(slide, emu(11.6), emu(7.04), emu(1.0), emu(0.36),
             [[(str(n), 9, False, GRAY)]], align=PP_ALIGN.RIGHT)


def _apply_bullet(p, level):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int((0.32 + level * 0.42) * IN)))
    pPr.set("indent", str(-int(0.32 * IN)))
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": "2E86C1" if level == 0 else "8A97A6"})
    buClr.append(srgb)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022" if level == 0 else "\u25AA"})
    pPr.append(buClr)
    pPr.append(buFont)
    pPr.append(buChar)


def bullets(slide, x, y, w, h, items):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (level, text, size, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6 if level else 10)
        p.space_before = Pt(2 if level else 5)
        p.line_spacing = 1.08
        _apply_bullet(p, level)
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def bullet_slide(prs, title, items, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    title_band(slide, title)
    bullets(slide, emu(0.85), emu(1.5), emu(11.6), emu(5.2), items)
    footer(slide, n)
    return slide


def place_image_fit(slide, path, x, y, max_w, max_h):
    img = Image.open(path)
    ar = img.width / img.height
    w = max_w
    h = w / ar
    if h > max_h:
        h = max_h
        w = h * ar
    x_c = x + (max_w - w) / 2
    slide.shapes.add_picture(str(path), emu(x_c), emu(y), width=emu(w), height=emu(h))
    return w, h


def image_slide(prs, title, path, caption, n, max_h=5.0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    title_band(slide, title)
    place_image_fit(slide, path, 0.55, 1.4, 12.2, max_h)
    add_text(slide, emu(0.7), emu(6.55), emu(12), emu(0.4),
             [[(caption, 12, False, GRAY)]])
    footer(slide, n)
    return slide


# ----------------------------- deck -----------------------------
def main():
    data = json.loads(DATA.read_text())
    plots = make_plots(data)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    B, G = DARK, GRAY
    n = 0

    def next_n():
        nonlocal n
        n += 1
        return n

    # 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_rect(s, emu(0.9), emu(2.25), emu(2.2), emu(0.08), ACCENT)
    add_text(s, emu(0.9), emu(2.45), emu(11.6), emu(1.4),
             [[("X2 Box-Pickup Deployment Analysis", 38, True, WHITE)]])
    add_text(s, emu(0.92), emu(3.85), emu(11), emu(0.9),
             [[("Hardware run logs · desired vs measured trajectories · reward fixes",
                18, False, RGBColor(0xC6, 0xD2, 0xE2))]])
    add_rect(s, emu(0.92), emu(5.25), emu(4.4), emu(0.02), RGBColor(0x3A, 0x4A, 0x66))
    add_text(s, emu(0.92), emu(5.45), emu(11), emu(0.5),
             [[("19 logged runs  ·  v27 + v30 policies  ·  50 Hz encoder CSVs",
                14, False, ACCENT)]])
    next_n()

    # 2 objective / what we logged
    bullet_slide(prs, "Hardware data logged", [
        (0, "Each deploy writes a 50 Hz CSV: phase, frame, base IMU quat/gyro, and per-joint pos_meas / vel_meas / tgt.", 18, False, B),
        (0, "Desired trajectory = policy PD target (`__tgt`). Actual = measured encoder (`__pos_meas`).", 18, False, B),
        (0, "19 runs analyzed: 4× v27 hybrid (Jul 27), 15× v30 pickup/hybrid (Jul 29).", 18, False, B),
        (0, "Primary metrics: foot-edge (ankle-roll deviation), base roll σ, leg tracking RMSE.", 18, False, B),
    ], next_n())

    # 3 IRL failure modes
    bullet_slide(prs, "IRL failure modes (v30 deploy)", [
        (0, "Feet stand on outer edges; most of the sole off the floor during pickup and hold.", 18, False, B),
        (0, "Torso oscillates in roll to rebalance on that narrow contact patch.", 18, False, B),
        (0, "Stepping / stance widen during pickup, hold, and set-down (slip penalty alone was insufficient).", 18, False, B),
        (0, "Backward fall risk after set-down when post-release upright hold was out of distribution.", 18, False, B),
        (0, "Median foot-edge on v30 hardware: 8.2° (target < 3°). Worst run: 13.1°.", 18, False, B),
    ], next_n())

    # 4 summary table
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    title_band(s, "Aggregate hardware metrics")
    place_image_fit(s, plots["summary"], 1.2, 1.7, 10.8, 3.4)
    bullets(s, emu(0.85), emu(5.35), emu(11.6), emu(1.5), [
        (0, "v30 did not fix edge-standing vs v27 — same ~8° median ankle-roll deviation.", 16, False, B),
        (0, "Sim contact terms treated an edge-loaded foot as planted; hardware did not.", 16, False, B),
    ])
    footer(s, next_n())

    # 5 foot-edge bar
    image_slide(prs, "Foot-edge severity across all runs",
                plots["foot_edge"],
                "RMS |ankle-roll| during pickup. Red = v30, gray = v27. Dashed = 3° flat-foot target.",
                next_n(), max_h=4.8)

    # 6 worst ankle
    image_slide(prs, "Worst run — edge-standing signature",
                plots["worst_ankle"],
                "Run 17:52:22. Commanded ankle-roll nearly flat; measured swings hard onto the edge.",
                next_n(), max_h=4.6)

    # 7 hybrid ankles/knees
    image_slide(prs, "Desired vs actual — ankles and knees",
                plots["hybrid_joints"],
                "Hybrid run 17:37:41. Shaded bands = ramp / settle / pickup / carry / setdown / done.",
                next_n(), max_h=5.0)

    # 8 hips
    image_slide(prs, "Desired vs actual — hips and ankle-pitch",
                plots["hybrid_hips"],
                "Same hybrid run. Hip tracking holds better than ankle-roll — failure is concentrated at the foot.",
                next_n(), max_h=5.0)

    # 9 base attitude
    image_slide(prs, "Base attitude — fixation wobble",
                plots["hybrid_base"],
                "Repeated roll swings during pickup = rebalancing on foot edges, not a clean plant.",
                next_n(), max_h=4.6)

    # 10 per-joint RMSE
    image_slide(prs, "Per-joint tracking error (legs)",
                plots["per_joint"],
                "Ankle-roll joints (red) dominate leg RMSE — matches the edge-standing observation.",
                next_n(), max_h=4.6)

    # 11 why old rewards missed it
    bullet_slide(prs, "Why prior rewards missed the failure", [
        (0, "`foot_slip`: zero when the foot is tilted but stationary — edge-stand is free.", 18, False, B),
        (0, "`feet_contact_loss`: edge still registers force in PhysX, so it counts as planted.", 18, False, B),
        (0, "`feet_anchor` / slip: address XY drift and sliding, not foot roll about the edge.", 18, False, B),
        (0, "No term measured ankle-roll link tilt → foot-tilt bounced 4°–29° across checkpoints with no trend.", 18, False, B),
    ], next_n())

    # 12 penalties added — table
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    title_band(s, "Penalties added to fix hardware failures")
    rows = [
        ("Term", "W", "What it fixes"),
        ("foot_slip", "−3", "Skating under load (sim → thrash on HW)"),
        ("feet_contact_loss", "−2", "Stepping loophole (lift to evade slip)"),
        ("feet_anchor", "−2→−1", "Stance widen / permanent XY drift"),
        ("foot_not_flat", "−3", "Edge-standing (sin of foot tilt, contact-gated)"),
        ("feet_edge_contact", "−2", "Edge counts as not planted (θ≳10°)"),
        ("ori tracking", "×2", "Left torso lean during pickup"),
        ("end-hold in motion", "—", "Fall after set-down (3 s upright practice)"),
    ]
    gt = s.shapes.add_table(len(rows), 3, emu(0.7), emu(1.45), emu(11.9), emu(5.2)).table
    gt.columns[0].width = emu(3.2)
    gt.columns[1].width = emu(1.4)
    gt.columns[2].width = emu(7.3)
    for r in range(len(rows)):
        for c in range(3):
            cell = gt.cell(r, c)
            cell.text = rows[r][c]
            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT
            para.font.size = Pt(14 if r else 15)
            para.font.bold = (r == 0) or (c == 0)
            if r == 0:
                para.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                para.font.color.rgb = DARK if c else ACCENT
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = emu(0.12)
    footer(s, next_n())

    # 13 equations for key new terms
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    title_band(s, "New foot terms (v31)")
    rows_eq = [
        ("Foot not flat — smooth tilt cost while loaded", "eq_not_flat"),
        ("Edge contact — hard count of edge-planted feet", "eq_edge"),
        ("Contact loss — airborne foot (stepping)", "eq_contact_loss"),
        ("Anchor — XY drift from reset plant", "eq_anchor"),
    ]
    y = 1.45
    for label, key in rows_eq:
        add_rect(s, emu(0.85), emu(y + 0.08), emu(0.09), emu(0.95), ACCENT)
        add_text(s, emu(1.1), emu(y), emu(5.2), emu(1.1),
                 [[(label, 14, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
        place_image_fit(s, plots[key], 6.3, y + 0.15, 6.4, 0.9)
        y += 1.2
    footer(s, next_n())

    # 14 what each term targets operationally
    bullet_slide(prs, "Mapping: symptom → penalty", [
        (0, "Edge-standing / oscillating on feet → `foot_not_flat` (−3) + `feet_edge_contact` (−2).", 18, False, B),
        (0, "Stepping during pickup/hold/set-down → `feet_contact_loss` (−2) + stronger `foot_slip` (−3).", 18, False, B),
        (0, "Stance growing apart → `feet_anchor` (drift from plant pose).", 18, False, B),
        (0, "Torso roll left on lift → double `motion_global_ref_orientation` weight.", 18, False, B),
        (0, "Fall after drop → bake 3 s end-hold into reference motion variants.", 18, False, B),
    ], next_n())

    # 15 hold / export caveats
    bullet_slide(prs, "Deployment notes from the data", [
        (0, "Hybrid deploy freezes the WBT clock at hold_frame (~261) for the walk splice — looks like a mid-hold stop.", 18, False, B),
        (0, "Full continuous pickup→set-down uses `deploy_x2_box_pickup.py`, not the hybrid script.", 18, False, B),
        (0, "Always pass `--box-policy` explicitly; default path previously loaded stale v27.", 18, False, B),
        (0, "Judge new checkpoints by foot-tilt + post-set-down upright, not mean reward alone.", 18, False, B),
    ], next_n())

    # 16 takeaways
    bullet_slide(prs, "Takeaways", [
        (0, "Hardware logs localize the failure to ankle-roll / foot tilt, not a generic tracking miss.", 18, False, B),
        (0, "Sim contact force ≠ flat sole — need an explicit flatness / edge objective.", 18, False, B),
        (0, "Planted-feet suite + foot-flat terms close the loopholes that produced IRL edge-standing and steps.", 18, False, B),
        (0, "Hold and set-down need dedicated sampling; short episodes never reach those phases.", 18, False, B),
    ], next_n())

    out = OUT_DIR / "X2_Box_Pickup_Deployment.pptx"
    prs.save(str(out))
    print(f"saved {out}  ({len(list(prs.slides))} slides)")
    print(f"assets in {ASSET_DIR}")


if __name__ == "__main__":
    main()
